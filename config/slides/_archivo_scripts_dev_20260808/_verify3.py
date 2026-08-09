# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pathlib import Path
from pptx import Presentation

root = Path(r"G:\Mi unidad\Trabajos\Empleo\UNIAJ\Cursos")

def slide_text(pptx, n):
    prs = Presentation(str(pptx))
    s = prs.slides[n-1]
    parts = []
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t: parts.append(t)
    return "\n".join(parts)

for curso, fname in [
    ("Programacion II", "Presentacion del Curso - Programacion II.pptx"),
    ("Arquitectura de Sistemas Computacionales", "Presentacion del Curso - Arquitectura de Sistemas Computacionales.pptx"),
]:
    p = root/curso/"Clases"/fname
    print("====", curso, "CONTENIDO")
    print(slide_text(p, 9)[:500])
    print()

# Arq clase1 mapa
p = list((root/"Arquitectura de Sistemas Computacionales"/"Clases").glob("Clase 1*/Presentacion.pptx"))[0]
prs = Presentation(str(p))
print("==== ARQ Clase1 slides with Mapa/Diagnóstico")
for i, s in enumerate(prs.slides, 1):
    texts=[]
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            texts.append(sh.text_frame.text.strip())
    j=" | ".join(texts)
    if "Mapa del bloque" in j or (i==4):
        print(f"slide {i}: {j[:300]}")
        # count shapes for visual richness
        print(f"  shapes={len(s.shapes)}")

# Seminario clase 1
sem = root/"Seminario de Sistemas"/"Clases"
print("==== Seminario Clases")
for d in sorted(sem.iterdir()):
    print(" ", d.name)