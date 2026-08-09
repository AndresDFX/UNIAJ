# -*- coding: utf-8 -*-
from pathlib import Path
from pptx import Presentation

root = Path(r"G:\Mi unidad\Trabajos\Empleo\UNIAJ\Cursos")

def texts(pptx):
    prs = Presentation(str(pptx))
    out = []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    out.append((i, t.replace("\n", " | ")[:160]))
    return out

# Verify Presentacion del Curso CONTENIDO
for curso, fname in [
    ("Programacion II", "Presentacion del Curso - Programacion II.pptx"),
    ("Seminario de Sistemas", "Presentacion del Curso - Seminario de Sistemas.pptx"),
    ("Bases de Datos II", "Presentacion del Curso - Bases de Datos II.pptx"),
    ("Arquitectura de Sistemas Computacionales", "Presentacion del Curso - Arquitectura de Sistemas Computacionales.pptx"),
]:
    p = root / curso / "Clases" / fname
    print("====", curso)
    for i, t in texts(p):
        if "Sesión" in t or "Sesion" in t or "Diagnóstico" in t or "Diagnostico" in t or "CONTENIDO" in t or "Clase 1" in t:
            if any(k in t for k in ["Sesión 0", "Sesion 0", "Diagnóstico ·", "Diagnostico ·", "Presentación del curso ·", "CONTENIDO", "Cómo trabajamos", "Como trabajamos"]):
                print(f"  slide {i}: {t[:140]}")

# Timeline slide in Arq Clase 1
p = list((root/"Arquitectura de Sistemas Computacionales"/"Clases").glob("Clase 1*/Presentacion.pptx"))[0]
print("==== ARQ Clase1 timeline", p)
for i, t in texts(p):
    if "Mapa" in t or "Diagnóstico" in t or "Diagnostico" in t or "0-10" in t or "Teoría" in t or "Teoria" in t:
        print(f"  slide {i}: {t[:140]}")

# Prog2 clase1
p = root/"Programacion II"/"Clases"/"Clase 1 - Introduccion a POO"/"Presentacion.pptx"
print("==== PROG2 Clase1")
for i, t in texts(p):
    if "Mapa" in t or "Diagnóstico" in t or "0-10" in t:
        print(f"  slide {i}: {t[:140]}")

# Seminario Clase 1 folder
sem = root/"Seminario de Sistemas"/"Clases"
print("==== Seminario Clases dirs")
for d in sorted(sem.iterdir()):
    if d.is_dir() and "Clase 1" in d.name:
        print(" ", d.name, list(d.glob("*"))[:5])