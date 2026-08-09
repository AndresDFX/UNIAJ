# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pathlib import Path
from pptx import Presentation

root = Path(r"G:\Mi unidad\Trabajos\Empleo\UNIAJ\Cursos")

def dump_contenido(pptx):
    prs = Presentation(str(pptx))
    for i, s in enumerate(prs.slides, 1):
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
        joined = "\n".join(texts)
        if "CONTENIDO" in joined and ("Clase 1" in joined or "Sesión" in joined or "Sesion" in joined):
            print(f"--- slide {i} ---")
            print(joined[:1200])
            print()

for curso, fname in [
    ("Programacion II", "Presentacion del Curso - Programacion II.pptx"),
    ("Arquitectura de Sistemas Computacionales", "Presentacion del Curso - Arquitectura de Sistemas Computacionales.pptx"),
    ("Bases de Datos II", "Presentacion del Curso - Bases de Datos II.pptx"),
    ("Seminario de Sistemas", "Presentacion del Curso - Seminario de Sistemas.pptx"),
]:
    print("====", curso)
    dump_contenido(root / curso / "Clases" / fname)

# check build sources for sesion0 item
for name in ["build_uniajc_prog2_curso.py", "build_uniajc_arq_curso.py"]:
    t = (root/".config/slides"/name).read_text(encoding="utf-8")
    print(name, "kind sesion0", 'kind": "sesion0"' in t or "kind': 'sesion0'" in t)
    # show contenido first items
    idx = t.find("contenido_clases_slides")
    print(t[idx:idx+500])