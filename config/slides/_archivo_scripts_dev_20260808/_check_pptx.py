from pptx import Presentation
from pptx.util import Inches, Emu
from pathlib import Path

files = [
    Path(r"Programacion II/Clases/Presentacion del Curso - Programacion II.pptx"),
    Path(r"Seminario de Sistemas/Clases/Presentacion del Curso - Seminario de Sistemas.pptx"),
    Path(r"Bases de Datos II/Clases/Presentacion del Curso - Bases de Datos II.pptx"),
    Path(r"Arquitectura de Sistemas Computacionales/Clases/Presentacion del Curso - Arquitectura de Sistemas Computacionales.pptx"),
]

def emu_to_in(e):
    return e / 914400

for f in files:
    prs = Presentation(str(f))
    found = None
    for i, s in enumerate(prs.slides):
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
        joined = " | ".join(texts)
        if "Herramientas del curso" in joined:
            pics = [sh for sh in s.shapes if sh.shape_type is not None and hasattr(sh, "image")]
            # picture shapes
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            pics = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            sizes = sorted([(emu_to_in(p.width), emu_to_in(p.height)) for p in pics], reverse=True)
            names = [t for t in texts if t and "Herramientas" not in t and "Gratis" not in t and len(t)<40]
            print(f.name)
            print("  slide", i+1, "pics", len(pics), "top sizes (in):", sizes[:6])
            print("  sample texts:", names[:8])
            found = True
            break
    if not found:
        print(f.name, "NO slide found")
