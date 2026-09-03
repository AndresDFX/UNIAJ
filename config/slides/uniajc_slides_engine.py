# -*- coding: utf-8 -*-
"""UNIAJC slides engine — reconstruido (patches + helpers)."""
from __future__ import annotations

import os

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import metrica_texto

NAVY = RGBColor(0x09, 0x52, 0x92)
CIAN = RGBColor(0x26, 0x9C, 0xCB)
AMARILLO = RGBColor(0xFF, 0xD0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x33, 0x33, 0x33)
SOFT = RGBColor(0x66, 0x66, 0x66)
ALT = RGBColor(0xF2, 0xF2, 0xF2)
INFO = RGBColor(0xE8, 0xF4, 0xFA)
ACLAR = RGBColor(0xFF, 0xF8, 0xE1)
WARN = RGBColor(0xFB, 0xE4, 0xE4)
RED = RGBColor(0xA0, 0x20, 0x30)
VERDE = RGBColor(0x1B, 0x7A, 0x4E)
LINK_TEAL = RGBColor(0x0D, 0x7A, 0x8C)
SW = 13.333
SH = 7.5
MARGIN = 0.5
CONTENT_W = SW - 2 * MARGIN

# ─────────────────────────────────────────────────────── cuerpo de las diapositivas
#
#: Tamano OBJETIVO del cuerpo de una diapositiva de contenido, en puntos.
#:
#: Los cinco cursos venian con el cuerpo escrito llamada por llamada: `size=13` en las
#: diapositivas densas, `size=16` en las sueltas, y unos 200 numeros distintos repartidos
#: por los builders. Esos numeros no eran una decision de diseno, eran ajuste manual: el
#: autor bajaba el tamano hasta que el texto entrara en la caja. Se leen bien en el monitor
#: de quien arma el deck y mal en una clase virtual, donde el estudiante recibe una ventana
#: compartida y recomprimida, a veces en un telefono.
#:
#: Ahora el tamano lo decide `bullets()`: parte de aqui y baja solo si el texto no cabe
#: (ver `CUERPO_MINIMO`). Para subir el cuerpo de TODOS los cursos se cambia este numero.
CUERPO_PT = 20

#: Suelo del ajuste. Por debajo no se baja: si el texto no cabe ni a 15 pt, lo que sobra es
#: texto en la diapositiva y no falta tamano de letra. Encogerlo mas taparia el problema;
#: `verificar_desborde.py` lo denuncia como choque para que se corrija el contenido.
CUERPO_MINIMO = 15

#: El bloque de codigo (`pseudo_code_slide`) va en su propia escala: son lineas que no se
#: pueden partir sin cambiar lo que dicen, en la caja mas apretada del motor. 14 pt es lo que
#: usaba fijo; el suelo es 11 porque por debajo un nombre de variable proyectado no se lee.
CODIGO_PT = 14
CODIGO_MINIMO = 11

ASSETS = Path(__file__).resolve().parent / "assets"
BRAND_PATH = Path(__file__).resolve().parents[1] / "universidades" / "uniajc.json"
QR_PADLET = str(ASSETS / "qr_padlet_uniajc.png")
OUTLOOK = str(ASSETS / "outlook_icon.png")
GMAIL = str(ASSETS / "gmail_icon.png")
PADLET_URL_PLACEHOLDER = "https://padlet.com/andres_dfx/uniaj-l77e9uu16trgdvcp"
PADLET_CLEAR_NOTE = "⋯ → Clear posts → código → Delete (rutina docente; no proyectar en clase)."

def load_brand():
    if BRAND_PATH.exists():
        try:
            return json.loads(BRAND_PATH.read_text(encoding="utf-8-sig"))
        except Exception:
            return {}
    return {}

# ───────────────────────────────────────── modo generico de las diapositivas de clase
#
# Que problema resuelve: una diapositiva de clase no puede dar por hecho lo que el docente
# decide semana a semana. Los decks nombraban la plataforma de entrega con su URL, anunciaban
# «taller domingo 23:59» y mencionaban quices, cuando el docente pone taller o no lo pone, y
# lo unico fijo de una sesion es EL TEMA. Proyectar un plazo que despues no se aplica es peor
# que no decir nada: el estudiante lo toma como compromiso.
#
# Donde SI se nombra la plataforma: en la Presentacion del Curso, que es el documento donde se
# explica una vez como se entrega. Por eso el modo generico esta ENCENDIDO por omision y son
# los cinco builders de `*_curso.py` los que lo apagan con `new_prs(generico=False)`. Al
# revés —apagado por omision— cualquier deck nuevo nacería nombrando la plataforma.
#
# Se aplica en `normalizar_inline` y en `_run`, que son los dos sitios por donde pasa TODO el
# texto que llega a una diapositiva. En `normalizar_inline` a proposito, y no solo al pintar:
# `bullets()` mide el texto para elegir el tamano, y tiene que medir el texto final.
GENERICO = True

_URL_PLATAFORMA = r"\s*\(\s*https?://[^)]*examlab[^)]*\)"

#: (patron, reemplazo), en orden: de lo mas especifico a lo mas general. El verificador
#: `verificar_generico.py` comprueba el resultado, asi que una regla que se quede corta se
#: detecta en vez de publicarse.
_REGLAS_GENERICO = [
    # 1. La herramienta del dia: lo que importa es el motor, no la plataforma que lo hospeda.
    (r"ExamLab\s*\(PostgreSQL[^)]*\)", "PostgreSQL en el navegador"),
    # 2. La URL y el modulo: nunca en una diapositiva.
    (_URL_PLATAFORMA, ""),
    # 3. El taller semanal, frase COMPLETA y antes que nada. Si la parte primero la regla
    #    del plazo queda «…esta semana: en la fecha que indique el docente», con «docente»
    #    dos veces y sin decir nada.
    (r"(?i)Taller de la semana en ExamLab[^.]*\.",
     "Si el docente asigna taller esta semana, él indica el canal y la fecha."),
    # 4. El plazo fijo. «(regla del Acuerdo)» y «cuando aplique» viajan pegados.
    (r"(?i)domingo 23:59\s*(\(regla del Acuerdo\))?\s*(cuando aplique(\s+el)?(\s+taller)?)?",
     "en la fecha que indique el docente"),
    (r"(?i)cuando aplique(\s+taller)?", "si el docente lo asigna"),
    (r"Entrega del taller en ExamLab", "Entrega del taller, si el docente lo asigna"),
    # 5. Los quices no se anuncian al estudiante.
    (r"\btaller y quiz\b", "taller"),
    (r"\bquices?\b", "las comprobaciones que el docente aplique"),
    # 6. Lo que quede del nombre propio.
    (r"\bExamLab\b", "la plataforma del curso"),
]


def texto_generico(t):
    """Quita de un texto de diapositiva lo que el docente decide cada semana.

    Idempotente: aplicarla dos veces da lo mismo, porque ninguna regla vuelve a introducir
    el patron de otra. Hace falta porque el texto pasa por `normalizar_inline` al medir y
    otra vez al pintar.
    """
    s = str(t)
    for pat, rep in _REGLAS_GENERICO:
        s = re.sub(pat, rep, s)
    # NO se colapsan los espacios repetidos: el prefijo de vineta que pone `bullets()` es
    # «–   » con tres espacios, y colapsarlo dejaba todas las vinetas de los decks de clase
    # con la sangria de un espacio mientras las de la Presentacion del Curso —que no pasa por
    # aqui— conservaban la de tres. Las reglas de arriba ya absorben el espacio que precede a
    # lo que borran, asi que no queda hueco doble que limpiar.
    #
    # La puntuacion se pega al texto solo cuando de verdad cierra la frase: sin el lookahead,
    # un «los .sql que pide el entregable» quedaba como «los.sql», porque la extension de
    # archivo empieza por punto y no es puntuacion.
    return re.sub(r"[ \t]+([.,;:])(?=\s|$)", r"\1", s)


def new_prs(generico=True):
    """Presentacion vacia 16:9. `generico=False` solo en la Presentacion del Curso.

    Es un interruptor global y no un atributo de la presentacion porque el texto se pinta
    desde `_run`/`_rich`, que reciben el parrafo y no la presentacion. Se fija al crear el
    deck, que es cuando se sabe de que tipo es, y cada builder construye un deck a la vez.
    """
    global GENERICO
    GENERICO = bool(generico)
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def rect(slide, l, t, w, h, color, line=False, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line_color or GRAY
    else:
        shape.line.fill.background()
    return shape


def rounded(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.12
    except Exception:
        pass
    return shape

def textbox(slide, l, t, w, h, anchor=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.auto_size = None
        try:
            tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.TOP: "t", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
        except Exception:
            pass
    return tf

def _run(run, text, size, color, bold=False, italic=False):
    run.text = texto_generico(text) if GENERICO else str(text)
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return run

def normalizar_inline(text):
    """Conversiones de marcado que hace `_rich` ANTES de pintar el texto.

    Existe aparte porque medir cuanto ocupa una vineta exige medir el texto que de verdad
    se va a pintar: `@@label@@` se convierte en negrita —que en Calibri es mas ancha— y un
    `code span` pasa a «comillas angulares», que cambian el numero de caracteres. Midiendo
    la cadena cruda, `bullets()` elegiria el tamano contra un texto que no es el final.
    """
    text = texto_generico(text) if GENERICO else str(text)
    text = re.sub(r"@@([^@]+)@@", r"**\1**", text)
    return re.sub(r"`([^`\n]+)`", r"«\1»", text)


def _rich(paragraph, text, size, color, bold=False, italic=False):
    """Soporta **negrita** y marcadores @@label@@ (negrita) en fragmentos.

    Los `code spans` de Markdown se convierten a «comillas angulares»: PowerPoint no
    tiene fuente monoespaciada dentro de una vineta, asi que los acentos graves salian
    impresos tal cual en la diapositiva («Los `GRANT` de hoy corren de verdad»). El
    fuente puede seguir escribiendose en Markdown, que es lo que necesita el guion .md.
    """
    if hasattr(paragraph, "clear"):
        try:
            paragraph.clear()
        except Exception:
            pass
    if paragraph.runs:
        paragraph.runs[0].text = ""
    text = normalizar_inline(text)
    # El contenido en negrita admite un `*` suelto: `[^*]+` no lo admitia, y una
    # vineta con «COUNT(*)» dentro del tramo en negrita no encontraba pareja, asi
    # que los dos asteriscos salian IMPRESOS en la diapositiva del estudiante y el
    # tramo resaltado se corria hasta el siguiente `**` del texto. Se encontro en la
    # Clase 6 de BD II, donde COUNT(*) contra COUNT(columna) es el nucleo de una
    # pregunta de 20 puntos. `\*(?!\*)` acepta el asterisco solitario y sigue sin
    # poder cruzar un cierre, y el `+?` toma el cierre mas cercano.
    parts = re.split(r"(\*\*(?:[^*]|\*(?!\*))+?\*\*)", text)
    first = True
    for part in parts:
        if not part:
            continue
        if first and paragraph.runs:
            run = paragraph.runs[0]
            first = False
        else:
            run = paragraph.add_run()
            first = False
        if part.startswith("**") and part.endswith("**"):
            _run(run, part[2:-2], size, color, bold=True, italic=italic)
        else:
            _run(run, part, size, color, bold=bold, italic=italic)
    return paragraph


def title_block(slide, title, sub=None):
    """Banda navy superior + título; retorna Y bajo el bloque."""
    bar_h = 1.05 if not sub else 1.28
    rect(slide, 0, 0, SW, bar_h, NAVY)
    rect(slide, 0, bar_h, SW, 0.045, AMARILLO)
    tf = textbox(slide, MARGIN, 0.18, CONTENT_W, 0.55 if not sub else 0.42)
    _run(tf.paragraphs[0].add_run(), title, 26, WHITE, bold=True)
    if sub:
        ts = textbox(slide, MARGIN, 0.62, CONTENT_W, 0.45)
        _rich(ts.paragraphs[0], sub, 14, RGBColor(0xB8, 0xD8, 0xEE))
    return bar_h + 0.18


ESPACIO_VINETA_PT = 8         # `space_after` entre viñetas


def ajustar(lineas, ancho, alto, objetivo, minimo, space_after_pt=0.0):
    """El tamaño más grande, de `objetivo` hacia abajo, con el que `lineas` cabe en la caja.

    Envoltorio de `metrica_texto.tamano_que_cabe` que normaliza el marcado inline primero:
    lo que se mide tiene que ser el texto que de verdad se va a pintar. Sin métrica
    disponible devuelve `objetivo`, que es el comportamiento de siempre.

    Los moldes que la usan son los tres que tenían texto pisando lo de abajo: las viñetas de
    contenido, el bloque de código y las columnas del mapa del bloque. No es cosmética: el
    motor no tiene autoajuste y PowerPoint dibuja el texto sobrante ENCIMA del pie o del
    borde, así que sin medir el deck se publica con líneas tapadas.
    """
    return metrica_texto.tamano_que_cabe(
        [normalizar_inline(x) for x in lineas], ancho, alto, objetivo, minimo,
        space_after_pt=space_after_pt)


def _con_vineta(item):
    """El texto de la viñeta tal como se va a pintar, con su guion si le falta."""
    raw = str(item)
    prefix = "" if raw.lstrip().startswith(("–", "-", "●", "•", "**1", "1.")) else "–   "
    return prefix + raw


def bullets(slide, items, top=1.4, size=None, width=None, left=None, minimo=None):
    """Lista de viñetas con soporte **negrita** y @@label@@, ajustada a su caja.

    `size=None` (lo normal) significa «el cuerpo del curso»: arranca en `CUERPO_PT` y baja
    de punto en punto hasta que el texto entre en la caja, sin pasar de `CUERPO_MINIMO`.
    Pasar un `size` explícito desactiva el ajuste y es lo que hacen los moldes con
    tipografía propia (código, notas, pasos numerados).

    Por qué se ajusta aquí y no a ojo en cada llamada: el motor no tiene autoajuste y
    PowerPoint dibuja el texto que no cabe **encima** de lo que haya debajo. Con el cuerpo
    en 20 pt eso le pasaba a decenas de diapositivas repartidas por los cinco cursos. El
    ajuste las deja todas dentro; las que ni al mínimo caben las denuncia
    `verificar_desborde.py`, porque ahí lo que sobra es texto, no tamaño de letra.
    """
    left = MARGIN if left is None else left
    width = CONTENT_W if width is None else width
    alto = max(0.8, SH - top - 0.5)
    lineas = [_con_vineta(x) for x in (items or [])]
    if size is None:
        size = metrica_texto.tamano_que_cabe(
            [normalizar_inline(x) for x in lineas],
            width, alto, CUERPO_PT, CUERPO_MINIMO if minimo is None else minimo,
            space_after_pt=ESPACIO_VINETA_PT)
    tf = textbox(slide, left, top, width, alto)
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(ESPACIO_VINETA_PT)
        _rich(p, linea, size, GRAY)
    return tf


def content_slide(prs, title, items, sub=None, idx=None, size=None):
    """Diapositiva de contenido con título UNIAJC + viñetas.

    `size=None` deja que `bullets()` elija el cuerpo (ver `CUERPO_PT`). Los builders ya no
    pasan el tamaño: los `size=13/14/15/16` que traían eran ajuste manual del autor para
    que el texto entrara, y ahora eso lo hace la medición.
    """
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    bullets(s, items, top=top + 0.05, size=size)
    footer_num(s, idx)
    return s


def _arrow(slide, x1, y1, x2, y2, color=None, dashed=False):
    """Conector recto con punta de flecha (triangle) al final."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color or SOFT
    conn.line.width = Pt(1.75)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(parse_xml(
            '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>'
        ))
    ln.append(parse_xml(
        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
    ))
    return conn


def diagram_boxes_slide(prs, title, boxes, arrows=None, sub=None, idx=None, note=None, legend=None):
    """Diagrama de cajas + flechas dibujado con formas reales (no imagen/placeholder).
    Util para C4 Context/Containers, diagramas de despliegue, ER simplificado, etc.

    boxes: [{id, label, x, y, w, h, color=CIAN, text_color=WHITE, size=13, dashed=False}]
    arrows: [{src, dst, label=None, dashed=False}]  (conecta el borde mas cercano entre cajas)
    legend: texto corto opcional bajo el diagrama (ej. leyenda de colores).
    """
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    box_map = {}
    for b in boxes:
        color = b.get("color", CIAN)
        shp = rounded(s, b["x"], b["y"], b["w"], b["h"], color)
        if b.get("dashed"):
            shp.line.color.rgb = GRAY
            shp.line.width = Pt(1)
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), b["label"], b.get("size", 13), b.get("text_color", WHITE), bold=True)
        box_map[b["id"]] = b
    for a in (arrows or []):
        src, dst = box_map[a["src"]], box_map[a["dst"]]
        sx, sy = src["x"] + src["w"] / 2, src["y"] + src["h"] / 2
        dx, dy = dst["x"] + dst["w"] / 2, dst["y"] + dst["h"] / 2
        if abs(dx - sx) >= abs(dy - sy):
            x1 = src["x"] + src["w"] if dx > sx else src["x"]
            x2 = dst["x"] if dx > sx else dst["x"] + dst["w"]
            y1, y2 = sy, dy
        else:
            y1 = src["y"] + src["h"] if dy > sy else src["y"]
            y2 = dst["y"] if dy > sy else dst["y"] + dst["h"]
            x1, x2 = sx, dx
        _arrow(s, x1, y1, x2, y2, dashed=a.get("dashed", False))
        if a.get("label"):
            ly = (y1 + y2) / 2 - 0.22
            if abs(x2 - x1) < 0.05:
                # Flecha vertical: el rotulo centrado en el punto medio quedaba
                # ATRAVESADO por la propia linea («usa · HTTPS» partido por la mitad en
                # el C4 Containers de la Clase 4 de ARQ). Se corre al lado de la linea,
                # y al otro lado si no cabe contra el borde derecho.
                lx, alin = x1 + 0.12, PP_ALIGN.LEFT
                if lx + 1.8 > MARGIN + CONTENT_W:
                    lx, alin = x1 - 1.92, PP_ALIGN.RIGHT
            else:
                lx, alin = (x1 + x2) / 2 - 0.9, PP_ALIGN.CENTER
            tb = textbox(s, max(MARGIN, lx), max(top, ly), 1.8, 0.35)
            tb.paragraphs[0].alignment = alin
            _run(tb.paragraphs[0].add_run(), a["label"], 10, SOFT, italic=True)
    if legend:
        lb = textbox(s, MARGIN, SH - 1.05, CONTENT_W, 0.35)
        _rich(lb.paragraphs[0], legend, 10.5, SOFT, italic=True)
    if note:
        nb = textbox(s, MARGIN, SH - 0.65, CONTENT_W, 0.45)
        _rich(nb.paragraphs[0], note, 11, SOFT, italic=True)
    footer_num(s, idx)
    return s


def table_content(
    prs,
    title,
    headers,
    rows,
    note=None,
    col_w=None,
    fs_body=12,
    idx=None,
):
    """Tabla simple para evaluación / cronograma."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title)
    n_cols = max(1, len(headers or []))
    if not col_w:
        col_w = [CONTENT_W / n_cols] * n_cols
    table_w = sum(col_w)
    n_rows = 1 + len(rows or [])
    shape = s.shapes.add_table(
        n_rows, n_cols, Inches(MARGIN), Inches(top + 0.1), Inches(table_w), Inches(0.42 * n_rows)
    )
    table = shape.table
    for i, w in enumerate(col_w):
        table.columns[i].width = Inches(w)
    for j, h in enumerate(headers or []):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        _run(p.add_run(), str(h), 12, WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r_i, row in enumerate(rows or [], start=1):
        for j in range(n_cols):
            cell = table.cell(r_i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            val = row[j] if j < len(row) else ""
            _rich(p, str(val), fs_body, GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ALT if r_i % 2 == 0 else WHITE
    if note:
        # La nota va DEBAJO de la tabla, no clavada al pie: con 4-5 filas quedaba a
        # tres pulgadas de distancia y parecia de otra diapositiva. Se estima el alto
        # con holgura para celdas de dos lineas y se limita para no salirse del area.
        y_note = min(SH - 0.95, top + 0.1 + 0.42 + 0.55 * len(rows or []) + 0.2)
        tn = textbox(s, MARGIN, y_note, CONTENT_W, 0.5)
        _rich(tn.paragraphs[0], note, 11, SOFT, italic=True)
    footer_num(s, idx)
    return s


def evaluacion_cortes_slide(prs, title, cortes, note=None, idx=None, sub=None):
    """Slide visual de evaluacion: una tarjeta grande por corte (en vez de tabla
    de texto plano). cortes: [{corte, ventana, desglose (str o list[str]), pct}].
    El % queda como numero grande — es lo que un docente/estudiante escanea
    primero — con una franja de color distinta por corte para lectura rapida."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    accents = [CIAN, AMARILLO, NAVY]
    n = max(1, len(cortes))
    gap = 0.35
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = SH - top - 0.9
    card_top = top + 0.15
    for i, c in enumerate(cortes):
        x = MARGIN + i * (card_w + gap)
        accent = accents[i % len(accents)]
        text_on_accent = NAVY if accent is AMARILLO else WHITE
        rounded(s, x, card_top, card_w, card_h, ALT)
        rect(s, x, card_top, card_w, 0.55, accent)
        p_head = textbox(s, x + 0.15, card_top + 0.05, card_w - 0.3, 0.45).paragraphs[0]
        _run(p_head.add_run(), f"Corte {c.get('corte', i + 1)}", 15, text_on_accent, bold=True)
        pct_box = textbox(s, x + 0.15, card_top + 0.65, card_w - 0.3, 0.95)
        p_pct = pct_box.paragraphs[0]
        p_pct.alignment = PP_ALIGN.CENTER
        _run(p_pct.add_run(), str(c.get("pct", "")), 40, accent, bold=True)
        vb = textbox(s, x + 0.15, card_top + 1.55, card_w - 0.3, 0.4)
        _rich(vb.paragraphs[0], c.get("ventana", ""), 11, SOFT, italic=True)
        desglose = c.get("desglose") or []
        if isinstance(desglose, str):
            desglose = [d.strip() for d in desglose.split("·")]
        db = textbox(s, x + 0.15, card_top + 2.0, card_w - 0.3, card_h - 2.1)
        for k, item in enumerate(desglose):
            p = db.paragraphs[0] if k == 0 else db.add_paragraph()
            p.space_after = Pt(6)
            _rich(p, f"–   {item}", 12.5, GRAY)
    if note:
        tn = textbox(s, MARGIN, SH - 0.6, CONTENT_W, 0.45)
        _rich(tn.paragraphs[0], note, 10.5, SOFT, italic=True)
    footer_num(s, idx)
    return s


def box_note_slide(prs, title, notes, idx=None):
    """Cajas info / aclaración / advertencia. notes: [(kind, texto), ...]."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title)
    palette = {
        "info": (INFO, CIAN, "Info"),
        "aclaracion": (ACLAR, AMARILLO, "Aclaración"),
        "advertencia": (WARN, RED, "Atención"),
    }
    y = top + 0.15
    for kind, text in notes or []:
        bg, accent, label = palette.get(kind, (ALT, CIAN, kind.title() if kind else "Nota"))
        h = 1.15
        rounded(s, MARGIN, y, CONTENT_W, h, bg)
        rect(s, MARGIN, y, 0.14, h, accent)
        tf = textbox(s, MARGIN + 0.35, y + 0.15, CONTENT_W - 0.55, h - 0.25)
        _run(tf.paragraphs[0].add_run(), label, 12, accent, bold=True)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        _rich(p2, text, 14, GRAY)
        y += h + 0.18
        if y > SH - 0.8:
            break
    footer_num(s, idx)
    return s

def footer_num(slide, idx=None):
    if idx is None:
        return
    tf = textbox(slide, SW - 1.1, SH - 0.42, 0.7, 0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p.add_run(), str(idx), 11, SOFT)

def _resolve_asset(name: str) -> Path | None:
    cands = [
        ASSETS / name,
        ASSETS / "herramientas" / name,
        ASSETS / "arq" / name,
    ]
    for c in cands:
        if c.exists():
            return c
    return None

def add_logo(slide, width=2.0, corner="left-top", mt=0.3, mr=0.5, variant="color"):
    mapping = {
        "white": "logo_uniajc_blanco.png",
        "color": "logo_uniajc_color.png",
        "hor_azul": "logo_uniajc_hor_azul.png",
        "hor_blanco": "logo_uniajc_hor_blanco.png",
    }
    fname = mapping.get(variant, "logo_uniajc.png")
    path = _resolve_asset(fname) or _resolve_asset("logo_uniajc.png")
    if not path:
        return None
    left = mr if "left" in corner else SW - width - mr
    top = mt if "top" in corner else SH - width * 0.4 - mt
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))

def _herramienta_logo_path(name: str) -> Path | None:
    """Resuelve el logo de una herramienta. Acepta tanto un nombre de archivo ya
    con extension ("dbfiddle.png") como un nombre suelto ("DB Fiddle").

    El match final ignora separadores (_ - espacio) porque los assets se nombran
    de forma inconsistente: `dbfiddle.png` vs `oracle_livesql.png` vs
    `play_with_docker.png`. Sin esa normalizacion, "DB Fiddle" -> "db_fiddle" no
    encontraba `dbfiddle.png` y TODA la grilla caia al placeholder de iniciales.
    """
    if not name:
        return None
    raw = str(name).strip()
    exts = (".png", ".jpg", ".jpeg", ".svg")
    # 1) Ya viene con extension: probarlo tal cual.
    if raw.lower().endswith(exts):
        p = _resolve_asset(raw)
        if p:
            return p
        base = Path(raw).stem
    else:
        base = raw
    # 2) Probar variantes del nombre + extension.
    stem = base.lower().replace(" ", "_")
    for cand_stem in {stem, stem.replace("_", ""), stem.replace("_", "-")}:
        for ext in exts:
            p = _resolve_asset(f"{cand_stem}{ext}")
            if p:
                return p
    # 3) Fuzzy ignorando separadores en AMBOS lados.
    def _norm(v: str) -> str:
        return re.sub(r"[^a-z0-9]", "", v.lower())

    target = _norm(base)
    if target:
        for folder in (ASSETS / "herramientas", ASSETS):
            if not folder.exists():
                continue
            for p in folder.iterdir():
                if not p.is_file() or p.suffix.lower() not in exts:
                    continue
                cand = _norm(p.stem)
                if cand and (cand == target or target in cand or cand in target):
                    return p
    return None

def course_cover(
    prs,
    materia,
    subtitulo,
    meta_lines,
    inicio_clase=None,
    hora_inicio_efectiva=None,
    pie_inicio=None,
    hora_arranque=None,
):
    """Portada del curso: banda azul institucional + filete amarillo + logo blanco.

    Pie opcional (solo portada Presentación del Curso): hora de arranque = horario
    oficial + 10 min. Alias: inicio_clase / hora_inicio_efectiva / pie_inicio /
    hora_arranque (p. ej. \"20:10\" o \"10:10\"). Wording: «Inicio de clase: HH:MM».
    No usar en PPTX de Clase N ni inventar otros pies.
    """
    s = blank(prs); bg_white(s)
    rect(s, 0, 0, SW, 3.2, NAVY)
    rect(s, 0, 3.2, SW, 0.08, AMARILLO)
    tf = textbox(s, MARGIN, 1.0, CONTENT_W, 1.6, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p.add_run(), materia, 34, WHITE, bold=True)
    if subtitulo:
        ps = tf.add_paragraph(); ps.alignment = PP_ALIGN.CENTER; ps.space_before = Pt(10)
        _run(ps.add_run(), subtitulo, 18, RGBColor(0xB8, 0xD8, 0xEE))
    add_logo(s, width=2.1, corner="left-top", mt=0.35, mr=0.5, variant="white")
    hora = inicio_clase or hora_inicio_efectiva or pie_inicio or hora_arranque
    meta_h = 2.85 if hora else 3.2
    tm = textbox(s, MARGIN, 3.65, CONTENT_W, meta_h)
    for i, ln in enumerate(meta_lines or []):
        p = tm.paragraphs[0] if i == 0 else tm.add_paragraph()
        p.space_after = Pt(8); _rich(p, ln, 15, GRAY)
    if hora:
        # Pie de portada: inicio efectivo = horario oficial del curso + 10 min
        rect(s, MARGIN, SH - 0.72, CONTENT_W, 0.02, CIAN)
        tf_pie = textbox(s, MARGIN, SH - 0.62, CONTENT_W, 0.4, anchor=MSO_ANCHOR.MIDDLE)
        pp = tf_pie.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        _rich(pp, f"Inicio de clase: **{hora}**", 16, NAVY)
    return s

def class_cover(prs, titulo, subtitulo=None, clase_n=None, idx=1):
    """Portada de Clase N: hero institucional limpio (sin bloque meta inferior).

    Solo: logo UNIAJC + título del tema + subtítulo corto opcional
    (p. ej. «Diagnóstico · CloudLite») + badge «Clase N».
    No poner aquí PI / agenda 120 min / gratis+navegador — van en la 2ª slide.
    """
    s = blank(prs)
    # Fondo hero navy a pantalla completa (evita franja + vacío blanco)
    rect(s, 0, 0, SW, SH, NAVY)
    # Acento superior cian + filete amarillo de marca
    rect(s, 0, 0, SW, 0.10, CIAN)
    rect(s, 0, 0.10, SW, 0.055, AMARILLO)
    add_logo(s, width=2.2, corner="left-top", mt=0.42, mr=0.55, variant="white")
    if clase_n is not None:
        badge_w, badge_h = 1.85, 0.42
        bx = SW - MARGIN - badge_w
        by = 0.42
        rounded(s, bx, by, badge_w, badge_h, CIAN)
        tb = textbox(s, bx, by, badge_w, badge_h, anchor=MSO_ANCHOR.MIDDLE)
        pb = tb.paragraphs[0]
        pb.alignment = PP_ALIGN.CENTER
        _run(pb.add_run(), f"Clase {clase_n}", 14, WHITE, bold=True)
    # Título centrado en el eje visual (composición hero)
    title_top = 2.55 if subtitulo else 2.85
    title_h = 2.0 if subtitulo else 1.6
    # Tipografía: títulos largos un poco más chicos para no desbordar
    tlen = len(titulo or "")
    title_pt = 40 if tlen <= 36 else (34 if tlen <= 56 else 28)
    tf = textbox(s, MARGIN + 0.3, title_top, CONTENT_W - 0.6, title_h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p.add_run(), titulo, title_pt, WHITE, bold=True)
    if subtitulo:
        # Filete amarillo corto bajo el título
        line_w = min(3.2, CONTENT_W * 0.28)
        rect(s, (SW - line_w) / 2, title_top + title_h + 0.05, line_w, 0.055, AMARILLO)
        ts = textbox(s, MARGIN + 0.3, title_top + title_h + 0.22, CONTENT_W - 0.6, 0.7)
        ps = ts.paragraphs[0]
        ps.alignment = PP_ALIGN.CENTER
        _run(ps.add_run(), subtitulo, 18, RGBColor(0xB8, 0xD8, 0xEE))
    else:
        line_w = min(3.2, CONTENT_W * 0.28)
        rect(s, (SW - line_w) / 2, title_top + title_h + 0.12, line_w, 0.055, AMARILLO)
    # Nº de slide discreto (tono claro sobre navy)
    if idx is not None:
        tn = textbox(s, SW - 1.2, SH - 0.48, 0.6, 0.3)
        pn = tn.paragraphs[0]
        pn.alignment = PP_ALIGN.RIGHT
        _run(pn.add_run(), str(idx), 11, RGBColor(0x7E, 0xB8, 0xD4), bold=True)
    return s

def tutor_slide(prs, nombre, credenciales, correo, rol=None, idx=None):
    """Diapositiva del docente: nombre grande, bullets de perfil y correo.

    Sin subtítulo tipo «Docente · <curso>». Logo UNIAJC color abajo-izquierda.
    Icono de correo: Outlook (asset outlook_icon.png).
    """
    s = blank(prs); bg_white(s)
    title_block(s, nombre)
    tx = 2.8
    ncred = max(1, len(credenciales or []))
    tc_h = min(2.6, 0.42 * ncred + 0.35)
    tc = textbox(s, tx, 2.25, SW - tx - MARGIN, tc_h)
    for i, ln in enumerate(credenciales):
        p = tc.paragraphs[0] if i == 0 else tc.add_paragraph()
        p.space_after = Pt(8); _rich(p, ln, 17, GRAY)
    ey = min(5.15, 2.25 + tc_h + 0.15)
    mail_icon = OUTLOOK if os.path.exists(OUTLOOK) else (GMAIL if os.path.exists(GMAIL) else None)
    if mail_icon:
        s.shapes.add_picture(mail_icon, Inches(tx), Inches(ey), height=Inches(0.32))
        te = textbox(s, tx + 0.55, ey - 0.02, SW - tx - 0.55 - MARGIN, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    else:
        te = textbox(s, tx, ey - 0.02, SW - tx - MARGIN, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    _run(te.paragraphs[0].add_run(), correo, 16, LINK_TEAL, bold=True)
    add_logo(s, width=1.8, corner="left-bottom", mt=0.35, mr=0.5, variant="color")
    footer_num(s, idx)
    return s

def block_timeline_slide(prs, title, slots, sub=None, idx=None, nota=None):
    """Timeline del bloque de HOY (no mapa del curso). slots: [(minutos, etiqueta), ...] o dicts.

    ``nota``: pie del timeline. El valor por omision dice «120 min» porque los cuatro
    cursos virtuales tienen bloques de 120; Introduccion a la Ingenieria los tiene de
    **90**, y con el texto clavado en el codigo la diapositiva mentia sobre su propia
    duracion. Pasar ``nota`` para cualquier curso que no sea de 120 min.
    """
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    norm = []
    for sl in slots or []:
        if isinstance(sl, dict):
            norm.append((sl.get("t") or sl.get("min") or "", sl.get("label") or sl.get("texto") or ""))
        else:
            norm.append((sl[0], sl[1] if len(sl) > 1 else ""))
    n = len(norm)
    if not n:
        footer_num(s, idx)
        return s
    y_line = top + 1.15
    rect(s, MARGIN, y_line, CONTENT_W, 0.08, CIAN)
    slot_w = CONTENT_W / n
    for i, (t, label) in enumerate(norm):
        cx = MARGIN + slot_w * i + slot_w / 2
        # nodo
        rounded(s, cx - 0.16, y_line - 0.12, 0.32, 0.32, NAVY if i % 2 == 0 else CIAN)
        tt = textbox(s, MARGIN + slot_w * i, top + 0.25, slot_w, 0.55)
        p = tt.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), str(t), 13, NAVY, bold=True)
        # La etiqueta de la columna llevaba 2.8 de alto fijo y 13 pt fijos, en una columna de
        # menos de 2 pulgadas de ancho: con 6 tramos y textos largos se metia por debajo de la
        # nota inferior. Ahora la caja llega hasta donde de verdad hay sitio —justo encima de
        # la nota— y el tamano se ajusta si aun asi no cabe.
        alto_lbl = max(0.8, (SH - 1.05) - (y_line + 0.45))
        ancho_lbl = slot_w - 0.1
        tb = textbox(s, MARGIN + slot_w * i + 0.05, y_line + 0.45, ancho_lbl, alto_lbl)
        pb = tb.paragraphs[0]
        pb.alignment = PP_ALIGN.CENTER
        _rich(pb, label, ajustar([label], ancho_lbl, alto_lbl, 13, 11), GRAY)
    # Nota inferior. Llevaba 0.4 de alto y 12 pt fijos, pensada para una linea; las notas
    # que explican el reparto de una sesion ocupan tres y se metian sobre el numero de
    # pagina. La caja llega ahora hasta justo encima del pie y el tamano se ajusta.
    txt_nota = nota or "Solo el bloque de **esta** clase · 120 min"
    y_nota = SH - 1.15
    alto_nota = (SH - 0.45) - y_nota
    tn = textbox(s, MARGIN, y_nota, CONTENT_W, alto_nota)
    tn.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(tn.paragraphs[0], txt_nota,
          ajustar([txt_nota], CONTENT_W, alto_nota, 12, 10), SOFT, italic=True)
    footer_num(s, idx)
    return s

def herramientas_slide(
    prs,
    items,
    title="Herramientas del curso",
    sub=None,
    idx=None,
    columns=None,
):
    """Grid de herramientas con logo PNG + nombre (Presentación del Curso).

    items: lista de dicts ``{name, logo, note?}`` o tuplas ``(name, logo[, note])``.
    ``logo`` = nombre de archivo en ``assets/herramientas/`` (p. ej. ``drawio.png``).
    Colocar **antes** de ``closing_slide``. Marca UNIAJC: bandas navy/cian + acento amarillo.
    """
    norm = []
    for raw in items or []:
        if isinstance(raw, dict):
            norm.append({
                "name": raw.get("name") or raw.get("nombre") or "",
                "logo": raw.get("logo") or raw.get("imagen"),
                "note": raw.get("note") or raw.get("caption") or raw.get("uso") or "",
            })
        elif isinstance(raw, (list, tuple)):
            norm.append({
                "name": raw[0] if len(raw) > 0 else "",
                "logo": raw[1] if len(raw) > 1 else None,
                "note": raw[2] if len(raw) > 2 else "",
            })
        else:
            norm.append({"name": str(raw), "logo": None, "note": ""})

    n = len(norm)
    if columns is None:
        if n <= 3:
            columns = n or 1
        elif n <= 6:
            columns = 3
        elif n <= 8:
            columns = 4
        else:
            columns = 4
    columns = max(1, min(columns, n or 1))
    rows = max(1, (n + columns - 1) // columns)

    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    # Nota corta bajo el título
    if not sub:
        hint = textbox(s, MARGIN, top - 0.05, CONTENT_W, 0.32)
        p = hint.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _rich(p, "Gratis · navegador o free tier · sin software de pago obligatorio", 12, SOFT)

    area_top = top + 0.28
    area_h = SH - area_top - 0.55
    gap_x, gap_y = 0.22, 0.18
    card_w = (CONTENT_W - gap_x * (columns - 1)) / columns
    card_h = (area_h - gap_y * (rows - 1)) / rows
    # Limitar altura de tarjeta para que no se vea estirada con pocas filas
    card_h = min(card_h, 2.55)

    for i, it in enumerate(norm):
        r, c = divmod(i, columns)
        # Fila incompleta (p. ej. 5 herramientas en 3 columnas): se centra en vez de
        # dejar el hueco pegado a la derecha, que descuadra la grilla.
        en_fila = min(columns, n - r * columns)
        sangria = (columns - en_fila) * (card_w + gap_x) / 2
        x = MARGIN + sangria + c * (card_w + gap_x)
        y = area_top + r * (card_h + gap_y)
        rounded(s, x, y, card_w, card_h, ALT)
        # filete superior marca
        rect(s, x, y, card_w, 0.07, NAVY)
        rect(s, x, y + 0.07, card_w, 0.03, AMARILLO)

        logo_path = _herramienta_logo_path(it.get("logo"))
        # Reparto vertical de la tarjeta: antes el logo se anclaba a 0.22 del borde
        # superior con un tope fijo de 1.15", asi que con pocas filas la tarjeta
        # quedaba con medio palmo de vacio abajo y la grilla se veia desbalanceada.
        # Ahora se calcula el alto del bloque (logo + nombre + nota) y se centra,
        # dejando que el logo crezca cuando hay sitio.
        texto_h = 0.42 + (0.26 if it.get("note") else 0.0)
        tope_logo = 1.5 if rows <= 2 else 1.15
        logo_side = min(card_w - 0.7, card_h - texto_h - 0.5, tope_logo)
        logo_side = max(0.55, logo_side)
        bloque_h = logo_side + 0.10 + texto_h
        lx = x + (card_w - logo_side) / 2
        ly = y + 0.16 + max(0.0, (card_h - 0.16 - bloque_h) / 2)
        if logo_path:
            try:
                from PIL import Image as _PILImage
                iw, ih = _PILImage.open(logo_path).size
            except Exception:
                iw, ih = (1, 1)
            if iw >= ih:
                pw, ph = logo_side, logo_side * ih / iw if iw else logo_side
            else:
                ph, pw = logo_side, logo_side * iw / ih if ih else logo_side
            # str(): python-pptx no acepta Path en add_picture (necesita ruta/stream).
            s.shapes.add_picture(
                str(logo_path),
                Inches(lx + (logo_side - pw) / 2),
                Inches(ly + (logo_side - ph) / 2),
                width=Inches(pw),
                height=Inches(ph),
            )
        else:
            rounded(s, lx, ly, logo_side, logo_side, INFO)
            tf = textbox(s, lx, ly, logo_side, logo_side, anchor=MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            initials = "".join(w[0] for w in (it.get("name") or "?").split()[:2]).upper()
            _run(p.add_run(), initials or "?", 16, NAVY, bold=True)

        name_y = ly + logo_side + 0.10
        name_h = min(texto_h, max(0.35, y + card_h - name_y - 0.06))
        tf = textbox(s, x + 0.08, name_y, card_w - 0.16, name_h, anchor=MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), it.get("name") or "", 13, NAVY, bold=True)
        if it.get("note"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(2)
            _run(p2.add_run(), it["note"], 10, SOFT)

    footer_num(s, idx)
    return s

def _normalize_contenido_item(it):
    """Acepta dict {n, tema, fecha?, tag?} o tupla (n, tema[, tag[, fecha]])."""
    if isinstance(it, dict):
        return it
    if isinstance(it, (list, tuple)):
        n = it[0]
        tema = it[1] if len(it) > 1 else ""
        tag = it[2] if len(it) > 2 else None
        fecha = it[3] if len(it) > 3 else None
        return {"n": n, "tema": tema, "tag": tag, "fecha": fecha}
    return {"n": "?", "tema": str(it)}

def _contenido_line(raw):
    it = _normalize_contenido_item(raw)
    n = it.get("n", "?")
    tema = (it.get("tema") or "").strip()
    fecha = it.get("fecha")
    tag = it.get("tag")
    kind = (it.get("kind") or "").lower()
    label = it.get("label")
    if label:
        # 2026-2: «Sesión N» != «Clase N» (13 sesiones / 15 clases de material),
        # asi que el rotulo puede venir dado explicitamente.
        line = f"–   **{label} –** {tema}"
    elif kind == "sesion0" or str(n) == "0":
        line = f"–   **Sesión 0 –** {tema}"
    else:
        line = f"–   **Clase {n} –** {tema}"
    if fecha:
        line += f" ({fecha})"
    if tag:
        line += f" · {tag}"
    return line


def contenido_clases_slide(prs, items, title="CONTENIDO", sub=None, idx=None, size=16, columns=None):
    """Slide estilo CONTENIDO: barra gris superior + título a la derecha + lista limpia,
    en UNA sola pagina (se pone en 2 columnas automaticamente si hay muchos items en
    vez de partir en varias slides — un temario partido en 2 paginas es mas dificil
    de escanear que uno compacto en 2 columnas).

    Reemplaza tablas densas #|Fecha|Tipo|Tema en Presentación del Curso.
    Cada ítem se renderiza como: – **Clase N –** Tema (fecha) · tag
      - tag breve solo si aporta (p. ej. «Parcial 1», «Autónoma»)
      - fecha corta opcional entre paréntesis

    items: lista de dicts {n, tema, fecha?, tag?} o tuplas (n, tema[, tag[, fecha]]).
    title: «CONTENIDO» (preferido) o «CONTENIDO DEL CURSO» / «CRONOGRAMA».
    sub: subtítulo discreto (p. ej. «Clases 1–8»).
    columns: fuerza 1 o 2 columnas; por defecto 2 si hay mas de 9 items.
    """
    items = list(items)
    s = blank(prs)
    bg_white(s)
    bar_h = 1.05
    rect(s, 0, 0, SW, bar_h, ALT)          # gris claro institucional #F2F2F2
    rect(s, 0, bar_h, SW, 0.035, CIAN)     # filete acento
    tf = textbox(s, MARGIN, 0, CONTENT_W, bar_h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p.add_run(), title, 28, NAVY, bold=True)

    y = bar_h + 0.16
    if sub:
        ts = textbox(s, MARGIN, y, CONTENT_W, 0.32)
        ps = ts.paragraphs[0]
        _rich(ps, sub, 13, SOFT)
        y += 0.38

    n_items = len(items)
    if columns is None:
        columns = 2 if n_items > 9 else 1
    # Tamano/espaciado se reduce con mas items por columna para que TODO quepa
    # en una sola pagina (mejor una lista compacta que partida en varias slides).
    per_col = -(-n_items // columns)  # ceil
    eff_size = size if per_col <= 8 else (size - 1 if per_col <= 11 else size - 2)
    space_after = 9 if per_col <= 8 else (6 if per_col <= 11 else 4)

    if columns == 1:
        body = textbox(s, MARGIN + 0.1, y, CONTENT_W - 0.2, SH - y - 0.5)
        for i, raw in enumerate(items):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.space_after = Pt(space_after)
            _rich(p, _contenido_line(raw), eff_size, GRAY)
    else:
        col_w = (CONTENT_W - 0.4) / 2
        left, right = items[:per_col], items[per_col:]
        for col_i, col_items in enumerate([left, right]):
            x = MARGIN + 0.1 + col_i * (col_w + 0.2)
            body = textbox(s, x, y, col_w - 0.1, SH - y - 0.5)
            for i, raw in enumerate(col_items):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.space_after = Pt(space_after)
                _rich(p, _contenido_line(raw), eff_size, GRAY)
    footer_num(s, idx)
    return s

def contenido_clases_slides(prs, items, title="CONTENIDO", sub=None, split_at=24, idx_start=None, size=16):
    """Compat: SIEMPRE emite 1 sola slide (2 columnas si hace falta) salvo que
    haya mas de `split_at` items (caso extremo), donde recien ahi parte en 2."""
    items = list(items)
    slides = []
    if len(items) <= split_at:
        slides.append(contenido_clases_slide(
            prs, items, title=title, sub=sub, idx=idx_start, size=size,
        ))
        return slides
    a, b = items[:split_at], items[split_at:]
    n2 = _normalize_contenido_item(b[0]).get("n", split_at + 1)
    n3 = _normalize_contenido_item(b[-1]).get("n", len(items))
    idx1 = idx_start
    idx2 = (idx_start + 1) if idx_start is not None else None
    slides.append(contenido_clases_slide(
        prs, a, title=title, sub=sub, idx=idx1, size=size,
    ))
    slides.append(contenido_clases_slide(
        prs, b, title=title, sub=f"Clases {n2}–{n3}", idx=idx2, size=size,
    ))
    return slides


# ---------- Portadas ----------

def padlet_slide(prs, idx=None, url=None):
    """Rompe-hielo Padlet (Presentación del Curso). Mismo QR/URL para todos los cursos UNIAJC.

    Solo contenido estudiante (sin Clear posts / checklist docente).
    """
    brand = load_brand()
    s = blank(prs); bg_white(s)
    top = title_block(s, "Rompe-hielo · Padlet del grupo")
    url = url or (brand.get("padlet") or {}).get("url") or PADLET_URL_PLACEHOLDER
    left_w = 7.2
    y = top + 0.12
    rounded(s, MARGIN, y, left_w, 4.55, INFO)
    rect(s, MARGIN, y, 0.14, 4.55, CIAN)
    tf = textbox(s, MARGIN + 0.35, y + 0.2, left_w - 0.55, 4.2)
    lines = [
        ("Escanea el **QR** o abre el enlace del tablero colaborativo.", 15),
        ("Deja **1 nota** en cada columna:", 15),
        ("**1.** ¿Quién soy? (nombre + una habilidad o hobby)", 14),
        ("**2.** Expectativa del curso (1 frase)", 14),
        ("**3.** Pregunta al docente / al grupo", 14),
        ("5–8 min · sin instalar · solo navegador.", 14),
        ("Enlace: **" + url + "**", 13),
    ]
    for i, (txt, sz) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        _rich(p, ("–   " if i < 2 or i >= 5 else "●   ") + txt if i not in (2, 3, 4) else txt, sz, GRAY)
    # QR a la derecha
    qr_box_x = MARGIN + left_w + 0.25
    qr_box_w = CONTENT_W - left_w - 0.25
    rounded(s, qr_box_x, y, qr_box_w, 4.55, ALT)
    if os.path.exists(QR_PADLET):
        try:
            from PIL import Image as _PILImage
            iw, ih = _PILImage.open(QR_PADLET).size
        except Exception:
            iw, ih = (1, 1)
        side = min(qr_box_w - 0.5, 3.5)
        h = side * ih / iw if iw else side
        left = Inches(qr_box_x + (qr_box_w - side) / 2)
        top_img = Inches(y + 0.35)
        s.shapes.add_picture(QR_PADLET, left, top_img, width=Inches(side), height=Inches(h))
        cap = textbox(s, qr_box_x + 0.1, y + 4.05, qr_box_w - 0.2, 0.4)
        p = cap.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p.add_run(), "QR oficial Padlet UNIAJC", 11, SOFT, bold=True)
    else:
        tf = textbox(s, qr_box_x + 0.2, y + 1.8, qr_box_w - 0.4, 1.0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _rich(p, "QR pendiente: `qr_padlet_uniajc.png`", 14, RED, bold=True)
    footer_num(s, idx)
    return s

def cards_grid_slide(prs, title, cards, sub=None, columns=None, idx=None):
    """Grid de tarjetas concepto (título + cuerpo). cards: [(titulo, cuerpo), ...] o dicts."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    norm = []
    for c in cards or []:
        if isinstance(c, dict):
            norm.append((c.get("title") or c.get("titulo") or "", c.get("body") or c.get("texto") or ""))
        elif isinstance(c, (list, tuple)):
            norm.append((c[0] if len(c) > 0 else "", c[1] if len(c) > 1 else ""))
        else:
            norm.append((str(c), ""))
    n = len(norm)
    if columns is None:
        columns = 3 if n >= 3 else max(1, n)
    columns = max(1, min(columns, n or 1))
    rows = max(1, (n + columns - 1) // columns)
    area_top = top + 0.15
    area_h = SH - area_top - 0.55
    gap_x, gap_y = 0.2, 0.18
    card_w = (CONTENT_W - gap_x * (columns - 1)) / columns
    card_h = min((area_h - gap_y * (rows - 1)) / rows, 2.8)
    accents = (NAVY, CIAN, AMARILLO, NAVY, CIAN, AMARILLO)
    for i, (ct, body) in enumerate(norm):
        r, c = divmod(i, columns)
        x = MARGIN + c * (card_w + gap_x)
        y = area_top + r * (card_h + gap_y)
        rounded(s, x, y, card_w, card_h, ALT)
        rect(s, x, y, card_w, 0.1, accents[i % len(accents)])
        tf = textbox(s, x + 0.18, y + 0.25, card_w - 0.36, 0.45)
        _run(tf.paragraphs[0].add_run(), ct, 15, NAVY, bold=True)
        if body:
            tb = textbox(s, x + 0.18, y + 0.75, card_w - 0.36, card_h - 0.95)
            _rich(tb.paragraphs[0], body, 13, GRAY)
    footer_num(s, idx)
    return s

def checklist_slide(prs, title, items, sub=None, idx=None):
    """Checklist de entregable (casillas visuales)."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    y = top + 0.15
    for i, it in enumerate(items or []):
        h = 0.62
        rounded(s, MARGIN, y, CONTENT_W, h, ALT)
        # casilla
        rect(s, MARGIN + 0.2, y + 0.16, 0.3, 0.3, WHITE, line=True)
        rect(s, MARGIN + 0.2, y + 0.16, 0.08, 0.3, CIAN)
        tf = textbox(s, MARGIN + 0.7, y + 0.1, CONTENT_W - 1.0, h - 0.18, anchor=MSO_ANCHOR.MIDDLE)
        _rich(tf.paragraphs[0], it, 15, GRAY)
        y += h + 0.12
        if y > SH - 0.7:
            break
    footer_num(s, idx)
    return s

def closing_slide(prs, big, lines, accent=None):
    s = blank(prs); bg_white(s)
    rect(s, 0, 0, SW, 0.18, CIAN)
    add_logo(s, width=1.6, corner="right-bottom", mt=0.3, mr=0.5, variant="color")
    tf = textbox(s, MARGIN, 1.9, CONTENT_W, 1.2, anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _rich(tf.paragraphs[0], big, 34, NAVY, bold=True)
    tl = textbox(s, MARGIN, 3.3, CONTENT_W, 1.8)
    for i, ln in enumerate(lines):
        p = tl.paragraphs[0] if i == 0 else tl.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6); _rich(p, ln, 16, GRAY)
    if accent:
        ta = textbox(s, MARGIN, 5.2, CONTENT_W, 0.7); ta.paragraphs[0].alignment = PP_ALIGN.CENTER
        _rich(ta.paragraphs[0], accent, 20, CIAN, bold=True)
    return s


# ---------- Herramientas del curso (Presentación del Curso) ----------
HERRAMIENTASASSETS = os.path.join(ASSETS, "herramientas")

def hook_slide(prs, hook, lines=None, eyebrow="Gancho de hoy", idx=None):
    """Diapositiva de apertura provocadora (poco texto, alto contraste)."""
    s = blank(prs)
    bg_white(s)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, 0.18, SH, AMARILLO)
    add_logo(s, width=1.8, corner="right-top", mt=0.35, mr=0.45, variant="white")
    te = textbox(s, MARGIN + 0.25, 1.4, CONTENT_W - 0.5, 0.4)
    pe = te.paragraphs[0]
    pe.alignment = PP_ALIGN.CENTER
    _run(pe.add_run(), eyebrow, 14, CIAN, bold=True)
    th = textbox(s, MARGIN + 0.25, 2.0, CONTENT_W - 0.5, 2.2, anchor=MSO_ANCHOR.MIDDLE)
    ph = th.paragraphs[0]
    ph.alignment = PP_ALIGN.CENTER
    _rich(ph, hook, 28, WHITE, bold=True)
    if lines:
        tl = textbox(s, MARGIN + 0.4, 4.5, CONTENT_W - 0.8, 1.8)
        for i, ln in enumerate(lines):
            p = tl.paragraphs[0] if i == 0 else tl.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(6)
            _rich(p, ln, 15, CIAN)
    footer_num(s, idx)
    return s

def image_side_slide(prs, title, image_path, items, side="right", sub=None,
                     caption=None, idx=None, size=15):
    """Texto + imagen (izquierda o derecha). Ideal diagramas C4 / capas / labs."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    gap = 0.3
    img_w = CONTENT_W * 0.48
    txt_w = CONTENT_W - img_w - gap
    y = top + 0.1
    h = SH - y - 0.55
    img_abs = _resolve_asset(image_path)
    if side == "left":
        ix, tx = MARGIN, MARGIN + img_w + gap
    else:
        tx, ix = MARGIN, MARGIN + txt_w + gap
    rounded(s, ix, y, img_w, h, ALT)
    if img_abs:
        try:
            from PIL import Image as _PILImage
            iw, ih = _PILImage.open(img_abs).size
            max_w, max_h = img_w - 0.3, h - (0.55 if caption else 0.3)
            scale = min(max_w / (iw / 96), max_h / (ih / 96), max_w, max_h)
            # Usar pulgadas directas acotadas al frame
            pw = min(max_w, max_h * iw / ih if ih else max_w)
            ph = pw * ih / iw if iw else max_h
            if ph > max_h:
                ph = max_h
                pw = ph * iw / ih if ih else max_w
            s.shapes.add_picture(
                img_abs,
                Inches(ix + (img_w - pw) / 2),
                Inches(y + 0.15),
                width=Inches(pw),
                height=Inches(ph),
            )
        except Exception:
            tf = textbox(s, ix + 0.2, y + h / 2 - 0.3, img_w - 0.4, 0.6, anchor=MSO_ANCHOR.MIDDLE)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            _run(tf.paragraphs[0].add_run(), "[Diagrama]", 14, SOFT)
    else:
        tf = textbox(s, ix + 0.2, y + h / 2 - 0.3, img_w - 0.4, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(tf.paragraphs[0].add_run(), f"[Asset: {image_path}]", 12, SOFT)
    if caption:
        tc = textbox(s, ix + 0.15, y + h - 0.45, img_w - 0.3, 0.35)
        tc.paragraphs[0].alignment = PP_ALIGN.CENTER
        _rich(tc.paragraphs[0], caption, 11, SOFT, italic=True)
    bullets(s, items, top=y + 0.1, size=size, width=txt_w, left=tx)
    footer_num(s, idx)
    return s

def link_callout_slide(prs, title, headline, url, notes=None, idx=None):
    """Slide para remarcar UN enlace obligatorio: caja grande verde + URL en tipografía grande."""
    s = blank(prs); bg_white(s)
    top = title_block(s, title)
    y = top + 0.25
    rounded(s, MARGIN, y, CONTENT_W, 2.4, INFO)
    rect(s, MARGIN, y, 0.18, 2.4, VERDE)
    tf = textbox(s, MARGIN + 0.45, y + 0.25, CONTENT_W - 0.7, 0.7)
    _rich(tf.paragraphs[0], headline, 20, NAVY, bold=True)
    uf = textbox(s, MARGIN + 0.45, y + 1.05, CONTENT_W - 0.7, 0.9, anchor=MSO_ANCHOR.MIDDLE)
    p = uf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p.add_run(), url, 22, VERDE, bold=True)
    y2 = y + 2.65
    for i, note in enumerate(notes or []):
        p = textbox(s, MARGIN, y2 + i * 0.55, CONTENT_W, 0.5).paragraphs[0]
        _rich(p, note, 14, GRAY)
    footer_num(s, idx)
    return s

def _partir_codigo(lineas):
    """Parte una lista de lineas de codigo en dos mitades por un limite de bloque.

    Busca el corte mas cercano a la mitad en el que las llaves estan equilibradas, para que
    una clase o una funcion no quede cortada entre columnas. Si no hay ningun punto asi
    —codigo sin llaves— parte por la mitad exacta.
    """
    n = len(lineas)
    medio = (n + 1) // 2
    saldo, candidatos = 0, []
    for i, ln in enumerate(lineas):
        saldo += ln.count("{") - ln.count("}")
        if saldo == 0:
            candidatos.append(i + 1)          # se puede cortar DESPUES de esta linea
    if candidatos:
        corte = min(candidatos, key=lambda c: (abs(c - medio), c))
        if 0 < corte < n:
            return lineas[:corte], lineas[corte:]
    return lineas[:medio], lineas[medio:]


def pseudo_code_slide(prs, title, lines, sub=None, idx=None, caption=None):
    """Bloque estilo terminal / pseudo-código visual (ADR, YAML, flujo)."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    y = top + 0.15
    h = SH - y - (0.9 if caption else 0.55)
    rounded(s, MARGIN, y, CONTENT_W, h, RGBColor(0x1E, 0x2A, 0x38))
    rect(s, MARGIN, y, CONTENT_W, 0.35, NAVY)
    td = textbox(s, MARGIN + 0.25, y + 0.05, CONTENT_W - 0.5, 0.28, anchor=MSO_ANCHOR.MIDDLE)
    _run(td.paragraphs[0].add_run(), "● ● ●", 10, CIAN, bold=True)
    # El bloque de codigo tambien se ajusta: es la caja mas apretada del motor y la que mas
    # se pasaba. Un `classDiagram` de 29 lineas a 14 pt pedia 6.9 pulgadas en una caja de
    # 4.5, asi que un tercio del diagrama salia encima del pie o fuera de la diapositiva —y
    # lo que se caia eran las multiplicidades, que era justo lo que la diapositiva ensenaba.
    # El `space_after` baja con el tamano: 4 pt entre 29 lineas son 1.5 pulgadas solo de aire.
    #
    # Si ni al minimo cabe en una columna, se parte en DOS antes de seguir encogiendo: a
    # 9 pt un nombre de variable proyectado no se lee, y perder o encoger el codigo es peor
    # que leerlo en dos mitades. Se parte por un limite de bloque (llaves equilibradas) para
    # no cortar una clase por la mitad.
    lineas = [str(x) for x in (lines or [])]
    ancho_cod, alto_cod = CONTENT_W - 0.7, h - 0.7
    size = metrica_texto.tamano_que_cabe(lineas, ancho_cod, alto_cod,
                                         CODIGO_PT, CODIGO_MINIMO, space_after_pt=4)
    columnas = [lineas]
    cabe = metrica_texto.alto_parrafos(lineas, ancho_cod, size, space_after_pt=4) <= alto_cod
    if lineas and not cabe:
        izq, der = _partir_codigo(lineas)
        ancho_col = (ancho_cod - 0.3) / 2
        size2 = min(
            metrica_texto.tamano_que_cabe(izq, ancho_col, alto_cod, CODIGO_PT,
                                          CODIGO_MINIMO, space_after_pt=4),
            metrica_texto.tamano_que_cabe(der, ancho_col, alto_cod, CODIGO_PT,
                                          CODIGO_MINIMO, space_after_pt=4))
        if size2 > size:
            columnas, size, ancho_cod = [izq, der], size2, ancho_col
    espacio = 4 * size / CODIGO_PT
    for col, trozo in enumerate(columnas):
        x = MARGIN + 0.35 + col * (ancho_cod + 0.3)
        tf = textbox(s, x, y + 0.5, ancho_cod, alto_cod)
        for i, ln in enumerate(trozo):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(espacio)
            _run(p.add_run(), ln, size, WHITE)
    if caption:
        tc = textbox(s, MARGIN, SH - 0.85, CONTENT_W, 0.35)
        tc.paragraphs[0].alignment = PP_ALIGN.CENTER
        _rich(tc.paragraphs[0], caption, 12, SOFT, italic=True)
    footer_num(s, idx)
    return s

def steps_visual_slide(prs, title, steps, sub=None, idx=None):
    """Pasos numerados grandes (demo / taller). steps: [str] o [(titulo, detalle)]."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    norm = []
    for st in steps or []:
        if isinstance(st, (list, tuple)):
            norm.append((st[0], st[1] if len(st) > 1 else ""))
        else:
            norm.append((str(st), ""))
    n = len(norm)
    y = top + 0.12
    area_h = SH - y - 0.55
    # El tope de altura evita filas absurdas con 2-3 pasos, pero con 4 pasos cortos
    # dejaba palmo y medio de blanco al pie. Se sube el tope y, si aun sobra sitio,
    # se centra el bloque en vez de dejarlo colgando arriba.
    row_h = min(1.15, area_h / max(n, 1) - 0.08)
    bloque_h = n * row_h + (n - 1) * 0.1
    y += max(0.0, (area_h - bloque_h) / 2)
    for i, (head, detail) in enumerate(norm):
        yy = y + i * (row_h + 0.1)
        rounded(s, MARGIN, yy, CONTENT_W, row_h, ALT if i % 2 == 0 else INFO)
        # badge número
        rounded(s, MARGIN + 0.15, yy + (row_h - 0.48) / 2, 0.48, 0.48, NAVY)
        tn = textbox(s, MARGIN + 0.15, yy + (row_h - 0.48) / 2, 0.48, 0.48, anchor=MSO_ANCHOR.MIDDLE)
        tn.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(tn.paragraphs[0].add_run(), str(i + 1), 16, WHITE, bold=True)
        ancho_f, alto_f = CONTENT_W - 1.1, row_h - 0.18
        # El paso llevaba 15/13 pt fijos en una fila de alto calculado a partir del NUMERO de
        # pasos, no de lo largo que sean: un «Paso 1. Arme en Google Docs el guion
        # cronometrado de doce minutos con...» pedia cinco lineas en una caja de dos y se
        # metia en la fila siguiente. El detalle va dos puntos por debajo del titulo, como
        # antes, y las dos bajan juntas.
        s_head = 15
        while s_head > 11:
            alto = metrica_texto.alto_parrafos([normalizar_inline(f"**{head}**")],
                                               ancho_f, s_head)
            if detail:
                alto += metrica_texto.alto_parrafos([normalizar_inline(detail)],
                                                    ancho_f, s_head - 2) \
                        - metrica_texto.INSET_V + 2 / 72
            if alto <= alto_f:
                break
            s_head -= 1
        tf = textbox(s, MARGIN + 0.8, yy + 0.1, ancho_f, alto_f, anchor=MSO_ANCHOR.MIDDLE)
        _rich(tf.paragraphs[0], f"**{head}**", s_head, NAVY, bold=True)
        if detail:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(2)
            _rich(p2, detail, s_head - 2, GRAY)
    footer_num(s, idx)
    return s

def two_column_slide(prs, title, left_items, right_items, left_title=None,
                     right_title=None, sub=None, idx=None, size=15):
    """Dos columnas con viñetas (comparaciones, teoría + práctica)."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    gap = 0.25
    col_w = (CONTENT_W - gap) / 2
    y = top + 0.1
    for i, (items, col_title) in enumerate(((left_items, left_title), (right_items, right_title))):
        x = MARGIN + i * (col_w + gap)
        rounded(s, x, y, col_w, SH - y - 0.55, ALT if i == 0 else INFO)
        rect(s, x, y, col_w, 0.08, NAVY if i == 0 else CIAN)
        ty = y + 0.2
        if col_title:
            tf = textbox(s, x + 0.18, ty, col_w - 0.36, 0.4)
            _run(tf.paragraphs[0].add_run(), col_title, 16, NAVY, bold=True)
            ty += 0.45
        bullets(s, items, top=ty, size=size, width=col_w - 0.36, left=x + 0.18)
    footer_num(s, idx)
    return s

def before_after_slide(prs, title, before_title, before_items, after_title, after_items,
                       sub=None, idx=None, size=15):
    """Comparación Antes / Después con acentos de marca."""
    s = blank(prs)
    bg_white(s)
    top = title_block(s, title, sub)
    gap = 0.35
    col_w = (CONTENT_W - gap) / 2
    y = top + 0.15
    h = SH - y - 0.55
    # Antes
    rounded(s, MARGIN, y, col_w, h, WARN)
    rect(s, MARGIN, y, col_w, 0.55, RED)
    tf = textbox(s, MARGIN + 0.15, y + 0.1, col_w - 0.3, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tf.paragraphs[0].add_run(), before_title, 16, WHITE, bold=True)
    bullets(s, before_items, top=y + 0.7, size=size, width=col_w - 0.35, left=MARGIN + 0.18)
    # Después
    x2 = MARGIN + col_w + gap
    rounded(s, x2, y, col_w, h, INFO)
    rect(s, x2, y, col_w, 0.55, NAVY)
    tf2 = textbox(s, x2 + 0.15, y + 0.1, col_w - 0.3, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tf2.paragraphs[0].add_run(), after_title, 16, WHITE, bold=True)
    bullets(s, after_items, top=y + 0.7, size=size, width=col_w - 0.35, left=x2 + 0.18)
    footer_num(s, idx)
    return s

def _parse_slot_duration(t):
    """Extrae duración en minutos desde '0-10', '10–35', '0–15 min', etc. Fallback None."""
    if t is None:
        return None
    s = str(t).strip().replace("–", "-").replace("—", "-")
    m = re.search(r"(\d+)\s*-\s*(\d+)", s)
    if m:
        a, b = int(m.group(1)), int(m.group(2))
        if b >= a:
            return max(1, b - a)
    m2 = re.search(r"(\d+)\s*min", s, re.I)
    if m2:
        return max(1, int(m2.group(1)))
    return None

