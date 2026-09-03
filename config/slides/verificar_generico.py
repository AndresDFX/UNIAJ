# -*- coding: utf-8 -*-
"""Comprueba que una diapositiva de CLASE no da por hecho lo que el docente decide.

La regla del curso
------------------
De una sesion lo unico fijo es **el tema**. Si hay taller, si hay una comprobacion escrita,
por que canal se entrega y con que plazo lo decide el docente esa semana. Una diapositiva que
anuncia «taller en <plataforma>, domingo 23:59» convierte en compromiso algo que a lo mejor no
se aplica, y el estudiante lo lee como tal.

Donde SI va el nombre de la plataforma: en la **Presentacion del Curso**, que es el documento
donde se explica una vez como se entrega y se puede corregir en un solo sitio.

Que comprueba
-------------
Sobre los `.pptx` de `<Curso>/Clases/`, separando por nombre de archivo los decks de clase de
las Presentaciones del Curso:

  - en un deck de CLASE no puede aparecer el nombre de la plataforma, su URL, un plazo fijo
    («domingo 23:59») ni la palabra «quiz»;
  - en la Presentacion del Curso si puede, y de hecho se avisa si NO aparece, porque entonces
    el estudiante no tiene donde enterarse.

El filtro que produce ese texto vive en `uniajc_slides_engine.texto_generico()` y se aplica
por deck con `new_prs(generico=...)`. Este verificador existe porque ese filtro es una lista
de reglas: si una se queda corta, es mejor que falle aqui que publicarse.

Uso
---
    python config/slides/verificar_generico.py
    python config/slides/verificar_generico.py "Programacion II/Clases/**/*.pptx"

Sale con codigo 1 si algo se cuela.
"""
from __future__ import annotations

import glob
import os
import re
import sys

from pptx import Presentation

#: Lo que NO puede aparecer en una diapositiva de clase. (etiqueta, patron).
PROHIBIDO_EN_CLASE = [
    ("nombre de la plataforma", r"ExamLab"),
    ("URL de la plataforma", r"https?://\S*examlab\S*"),
    ("plazo fijo", r"domingo\s+23:59"),
    ("quiz", r"\bqui(?:z|ces)\b"),
]

#: En la Presentacion del Curso se ESPERA encontrar esto: es el unico sitio donde se dice.
ESPERADO_EN_CURSO = ("ExamLab", r"ExamLab")

MARCA_CURSO = "presentacion del curso"


def _texto(ruta):
    prs = Presentation(ruta)
    out = []
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for pa in sh.text_frame.paragraphs:
                t = "".join(r.text for r in pa.runs).strip()
                if t:
                    out.append((i, t))
    return out


def _es_curso(ruta):
    return MARCA_CURSO in os.path.basename(ruta).lower()


def revisar(ruta):
    """`(fallos, avisos)` de un .pptx."""
    fallos, avisos = [], []
    lineas = _texto(ruta)
    if _es_curso(ruta):
        etiqueta, pat = ESPERADO_EN_CURSO
        if not any(re.search(pat, t, re.I) for _, t in lineas):
            avisos.append("no menciona %s: el estudiante no tiene donde enterarse de "
                          "como se entrega" % etiqueta)
        return fallos, avisos
    for etiqueta, pat in PROHIBIDO_EN_CLASE:
        for n, t in lineas:
            m = re.search(pat, t, re.I)
            if m:
                fallos.append("sl%d · %s: «%s»" % (n, etiqueta, t[:90]))
    return fallos, avisos


def main(argv=None):
    pats = argv or ["*/Clases/**/*.pptx"]
    rutas = sorted({r for p in pats for r in glob.glob(p, recursive=True)
                    if r.lower().endswith(".pptx")})
    if not rutas:
        print("No hay .pptx que casen con el patron.")
        return 1
    n_fallos = n_curso = 0
    for ruta in rutas:
        fallos, avisos = revisar(ruta)
        if _es_curso(ruta):
            n_curso += 1
        etq = (os.path.basename(ruta) if _es_curso(ruta)
               else os.path.basename(os.path.dirname(ruta)))
        if fallos:
            n_fallos += len(fallos)
            print("MAL  %s" % etq[:60])
            for f in fallos[:6]:
                print("       " + f)
            if len(fallos) > 6:
                print("       ... y %d mas" % (len(fallos) - 6))
        for a in avisos:
            print("AVISO %s: %s" % (etq[:50], a))
    print("\n  %d decks revisados (%d Presentacion del Curso) · %d fallos"
          % (len(rutas), n_curso, n_fallos))
    return 1 if n_fallos else 0


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    raise SystemExit(main(sys.argv[1:] or None))
