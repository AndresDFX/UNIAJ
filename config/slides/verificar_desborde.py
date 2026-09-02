# -*- coding: utf-8 -*-
"""Comprueba que el texto de un .pptx no se sale de su sitio.

Por que existe
--------------
Subir el cuerpo de las diapositivas de 13-16 pt a 20 pt es lo que las hace legibles en
una clase virtual, donde el estudiante ve una ventana compartida y recomprimida. Pero el
motor de slides NO tiene ajuste automatico: `bullets()` pinta el texto en una caja de alto
fijo y PowerPoint lo dibuja igual si no cabe, encima de lo que haya debajo. Un cuerpo mas
grande sin medir cambia «se lee mal» por «la ultima linea tapa el numero de pagina», que
es peor.

La medida vive en `metrica_texto` —la misma que usa `bullets()` para elegir el tamano, para
que el motor y este verificador no puedan discrepar—. Lo que decide aqui es que cuenta como
defecto, y eso tiene una regla que no es obvia:

**Pasarse del alto de la caja NO es un defecto por si solo.** Muchos moldes declaran
   cajas deliberadamente cortas —una fila de checklist mide 0.44 pulgadas— y centran una
linea dentro con `anchor=ctr`. El texto asoma de la caja y se ve perfecto. Lo que si es un
defecto es que el texto llegue a PISAR la forma de abajo o el borde de la diapositiva, y eso
es lo unico que se reporta.

Uso
---
    python config/slides/verificar_desborde.py "<Curso>/Clases/**/*.pptx"
    python config/slides/verificar_desborde.py <archivo.pptx> --cuerpo 20

`--cuerpo N` no modifica nada: proyecta que pasaria si el cuerpo fuera de N puntos, para
poder decidir ANTES de tocar un builder.

Sale con codigo 1 si algun texto pisa algo, para poder encadenarlo en una verificacion.
"""
from __future__ import annotations

import argparse
import glob
import os
import sys

from pptx import Presentation

import metrica_texto

EMU = 914400.0
SH_POR_DEFECTO = 7.5          # alto de la diapositiva 16:9 que usa el motor
MARGEN = 0.02                 # pulgadas de tolerancia: por debajo es ruido de medicion


def alto_texto(tf, ancho_in: float, forzar: float | None = None) -> float:
    """Alto en pulgadas que ocupa el texto de `tf` ajustado a `ancho_in`.

    Se apoya en `metrica_texto`, que es la misma medida que usa `bullets()` al elegir el
    tamano: si divergieran, el motor podria dar por bueno un cuerpo que este verificador
    denuncia. `forzar` sustituye el tamano de todos los runs, para proyectar otro cuerpo sin
    regenerar el archivo.
    """
    parrafos, tamanos, espacios = [], [], []
    for p in tf.paragraphs:
        tramos = [(r.text, bool(r.font.bold)) for r in p.runs if r.text]
        if not tramos:
            parrafos.append([("", False)])
            tamanos.append(forzar or 18)
        else:
            parrafos.append(tramos)
            tamanos.append(forzar or max((r.font.size.pt if r.font.size else 18)
                                         for r in p.runs if r.text))
        espacios.append(((p.space_after.pt if p.space_after else 0),
                         (p.space_before.pt if p.space_before else 0)))
    total = 0.0
    n = len(parrafos)
    for i, (par, size, (sa, sb)) in enumerate(zip(parrafos, tamanos, espacios)):
        # un parrafo a la vez: cada uno puede tener su propio tamano
        total += metrica_texto.alto_parrafos([par], ancho_in, size) - metrica_texto.INSET_V
        total += sb / 72
        if i < n - 1:
            total += sa / 72
    return total + metrica_texto.INSET_V


def _anchor(tf) -> str:
    try:
        return tf._txBody.bodyPr.get("anchor", "t") or "t"
    except Exception:
        return "t"


def _techo_debajo(slide, forma, alto_diapo: float) -> float:
    """Y (pulgadas) de lo primero que hay debajo de `forma` y se cruza con ella.

    Si no hay nada, el borde inferior de la diapositiva. Solo cuentan las formas que
    solapan en horizontal: una nota en la columna de al lado no estorba.
    """
    x0, x1 = forma.left / EMU, (forma.left + forma.width) / EMU
    fondo_caja = forma.top / EMU + forma.height / EMU
    techo = alto_diapo
    for otra in slide.shapes:
        if otra is forma:
            continue
        oy = otra.top / EMU
        if oy < fondo_caja - MARGEN:
            continue
        ox0, ox1 = otra.left / EMU, (otra.left + otra.width) / EMU
        if ox1 <= x0 + MARGEN or ox0 >= x1 - MARGEN:
            continue
        techo = min(techo, oy)
    return techo


def _cuerpo_de(slide, minimo: int):
    """La forma con mas texto de la diapositiva, o None. Titulos y pies quedan fuera."""
    mejor = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        n = sum(len(r.text) for p in sh.text_frame.paragraphs for r in p.runs)
        if n < minimo:
            continue
        if mejor is None or n > mejor[0]:
            mejor = (n, sh)
    return mejor


def _titulo_de(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.paragraphs:
            t = "".join(r.text for r in sh.text_frame.paragraphs[0].runs)
            if t and len(t) < 70:
                return t
    return ""


def revisar(ruta: str, forzar: float | None = None, minimo: int = 80):
    """`(choques, tamanos)` de un .pptx.

    `choques` es una lista de dicts con la diapositiva, su titulo y cuantas pulgadas pisa.
    `tamanos` es el conjunto de tamanos de cuerpo encontrados, para ver de un vistazo si
    un deck ya quedo en el cuerpo objetivo.
    """
    prs = Presentation(ruta)
    alto_diapo = (prs.slide_height / EMU) if prs.slide_height else SH_POR_DEFECTO
    choques, tamanos = [], set()
    for i, slide in enumerate(prs.slides, 1):
        mejor = _cuerpo_de(slide, minimo)
        if not mejor:
            continue
        _, sh = mejor
        tf = sh.text_frame
        tamanos |= {round(r.font.size.pt) for p in tf.paragraphs for r in p.runs
                    if r.font.size and r.text.strip()}
        necesita = alto_texto(tf, sh.width / EMU, forzar=forzar)
        top, alto_caja = sh.top / EMU, sh.height / EMU
        # con anclaje centrado el texto crece hacia los dos lados desde el centro
        fondo = (top + alto_caja / 2 + necesita / 2) if _anchor(tf) == "ctr" else top + necesita
        techo = _techo_debajo(slide, sh, alto_diapo)
        if fondo > techo + MARGEN:
            choques.append({"slide": i, "titulo": _titulo_de(slide),
                            "pisa": round(fondo - techo, 2),
                            "necesita": round(necesita, 2), "caja": round(alto_caja, 2)})
    return choques, tamanos


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("patron", nargs="+", help="archivo .pptx o patron glob (admite **)")
    ap.add_argument("--cuerpo", type=float, default=None,
                    help="proyecta el resultado con este cuerpo, sin modificar nada")
    ap.add_argument("-q", "--solo-choques", action="store_true",
                    help="imprime solo los decks con choques")
    a = ap.parse_args(argv)

    rutas = sorted({r for p in a.patron for r in glob.glob(p, recursive=True)
                    if r.lower().endswith(".pptx")})
    if not rutas:
        print("No hay .pptx que casen con el patron.")
        return 1

    total = 0
    for ruta in rutas:
        choques, tamanos = revisar(ruta, forzar=a.cuerpo)
        total += len(choques)
        if a.solo_choques and not choques:
            continue
        etiqueta = os.path.basename(os.path.dirname(ruta)) or os.path.basename(ruta)
        if etiqueta == "Clases":
            etiqueta = os.path.basename(ruta)
        print("%-52s cuerpo %-22s %s" % (
            etiqueta[:52],
            ",".join(str(t) for t in sorted(tamanos)) or "-",
            "OK" if not choques else "PISA en %s" % ", ".join(
                "sl%d (+%.2f\")" % (c["slide"], c["pisa"]) for c in choques)))
    print("\n  %d diapositivas con texto que pisa algo, en %d decks%s" % (
        total, len(rutas), " (proyectado a cuerpo %g)" % a.cuerpo if a.cuerpo else ""))
    return 1 if total else 0


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    raise SystemExit(main())
