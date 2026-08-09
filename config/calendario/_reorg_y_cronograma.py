# -*- coding: utf-8 -*-
"""Reorganiza Plan curso / Entregas docente + cronograma 1-15 en PPTX."""
from __future__ import annotations

import json
import re
import shutil
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]

CURSOS = {
    "programacion_ii": {
        "folder": "Programacion II",
        "build": "build_uniajc_prog2_curso.py",
        "csv_cfg": "eventos_programacion_ii_2026-2.csv",
        "rows": [
            ["1", "12/08", "Presencial", "Introducción a POO"],
            ["2", "19/08", "Virtual", "Colecciones dinámicas ArrayList"],
            ["3", "26/08", "Virtual", "Pilas y colas"],
            ["4", "02/09", "Virtual", "Mapas y conjuntos"],
            ["5", "09/09", "Presencial", "Interfaces gráficas GUI · Parcial 1"],
            ["6", "16/09", "Virtual", "Eventos y controladores"],
            ["7", "23/09", "Virtual", "Patrones de diseño"],
            ["8", "30/09", "Virtual", "Documentación y QA"],
            ["9", "07/10", "Virtual", "Refactorización con IA"],
            ["10", "14/10", "Presencial", "Persistencia de archivos · Parcial 2"],
            ["11", "21/10", "Virtual", "Revisión de código cruzada"],
            ["12", "28/10", "Virtual", "Integración de módulos"],
            ["13", "04/11", "Virtual", "Control de excepciones"],
            ["14", "11/11", "Virtual", "Preparación presentación final"],
            ["15", "18/11", "Presencial", "Evaluación de proyectos + cierre · Parcial 3"],
        ],
    },
    "seminario": {
        "folder": "Seminario de Sistemas",
        "build": "build_uniajc_seminario_curso.py",
        "csv_cfg": "eventos_seminario_2026-2.csv",
        "rows": [
            ["1", "13/08", "Presencial", "Acuerdo y conceptos"],
            ["2", "20/08", "Virtual", "Ciclos de vida"],
            ["3", "27/08", "Virtual", "Metodologías tradicionales"],
            ["4", "03/09", "Virtual", "Metodologías ágiles"],
            ["5", "10/09", "Presencial", "Caso estudio evaluativo · Parcial 1"],
            ["6", "17/09", "Virtual", "Requerimientos de software"],
            ["7", "24/09", "Virtual", "Historias de usuario"],
            ["8", "01/10", "Virtual", "Introducción a UML"],
            ["9", "08/10", "Virtual", "Casos de uso"],
            ["10", "15/10", "Presencial", "Caso estudio evaluativo · Parcial 2"],
            ["11", "22/10", "Virtual", "Avance proyecto integrador"],
            ["12", "29/10", "Virtual", "Diagramas UML avanzados"],
            ["13", "05/11", "Virtual", "Diseño de interfaces"],
            ["14", "12/11", "Virtual", "Evaluación final (prep. sustentación)"],
            ["15", "19/11", "Presencial", "Sustentación de proyectos + cierre · Parcial 3"],
        ],
    },
    "bases_datos_ii": {
        "folder": "Bases de Datos II",
        "build": "build_uniajc_bd2_curso.py",
        "csv_cfg": "eventos_bases_datos_ii_2026-2.csv",
        "rows": [
            ["1", "10/08", "Presencial", "Presentación · Revisión BD I"],
            ["2", "17/08", "Autónoma", "Administración de BD (festivo)"],
            ["3", "24/08", "Virtual", "Procedimientos almacenados"],
            ["4", "31/08", "Virtual", "Funciones y disparadores"],
            ["5", "07/09", "Presencial", "Seguridad y respaldo · Parcial 1"],
            ["6", "14/09", "Virtual", "Optimización de consultas"],
            ["7", "21/09", "Virtual", "Índices y particionamiento"],
            ["8", "28/09", "Virtual", "Tuning de bases de datos"],
            ["9", "05/10", "Presencial", "Gestión de transacciones · Parcial 2"],
            ["10", "12/10", "Autónoma", "Control de concurrencia (festivo)"],
            ["11", "19/10", "Virtual", "Avance del proyecto final"],
            ["12", "26/10", "Virtual", "Integración de aplicaciones externas"],
            ["13", "02/11", "Autónoma", "Análisis de casos reales (festivo)"],
            ["14", "09/11", "Presencial", "Prep. presentación final · Parcial 3"],
            ["15", "16/11", "Autónoma", "Presentación del proyecto + cierre (festivo)"],
        ],
    },
    "arquitectura": {
        "folder": "Arquitectura de Sistemas Computacionales",
        "build": "build_uniajc_arq_curso.py",
        "csv_cfg": "eventos_arquitectura_2026-2.csv",
        "rows": [
            ["1", "10/08", "Presencial", "Presentación · Intro arquitecturas cloud"],
            ["2", "17/08", "Autónoma", "IaaS / PaaS / SaaS (festivo)"],
            ["3", "24/08", "Virtual", "Virtualización y contenedores"],
            ["4", "31/08", "Virtual", "Microservicios"],
            ["5", "07/09", "Presencial", "Arquitecturas distribuidas · Parcial 1"],
            ["6", "14/09", "Virtual", "Seguridad en la nube"],
            ["7", "21/09", "Virtual", "Redes y almacenamiento cloud"],
            ["8", "28/09", "Virtual", "Monitoreo y optimización"],
            ["9", "05/10", "Presencial", "CI/CD · Parcial 2"],
            ["10", "12/10", "Autónoma", "Costos y sostenibilidad cloud (festivo)"],
            ["11", "19/10", "Virtual", "Avance del proyecto final"],
            ["12", "26/10", "Virtual", "Pruebas de rendimiento"],
            ["13", "02/11", "Autónoma", "Escalabilidad automática (festivo)"],
            ["14", "09/11", "Presencial", "Prep. presentación final · Parcial 3"],
            ["15", "16/11", "Autónoma", "Presentación del proyecto + cierre (festivo)"],
        ],
    },
}


def read_text(path: Path) -> str:
    raw = path.read_bytes()
    if raw.startswith((b"\xff\xfe", b"\xfe\xff")) or (len(raw) > 3 and raw[1] == 0):
        return raw.decode("utf-16")
    if raw.startswith(b"\xef\xbb\xbf"):
        return raw.decode("utf-8-sig")
    return raw.decode("utf-8")


def write_text(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8", newline="\n")


def move_file(src: Path, dst: Path) -> str:
    if not src.exists():
        return f"SKIP missing {src}"
    dst.parent.mkdir(parents=True, exist_ok=True)
    if dst.exists():
        if src.resolve() == dst.resolve():
            return f"OK already {dst.relative_to(ROOT)}"
        dst.unlink()
    shutil.move(str(src), str(dst))
    return f"MOVED {src.relative_to(ROOT)} -> {dst.relative_to(ROOT)}"


def reorganize() -> list[str]:
    log: list[str] = []
    for meta in CURSOS.values():
        folder = ROOT / meta["folder"]
        plan = folder / "Plan curso"
        entregas = folder / "Entregas docente"
        plan.mkdir(exist_ok=True)
        entregas.mkdir(exist_ok=True)

        # CSV eventos -> Plan curso
        log.append(
            move_file(
                entregas / "calendario_eventos_2026-2.csv",
                plan / "calendario_eventos_2026-2.csv",
            )
        )
        # CALENDARIO md root -> Plan curso
        log.append(
            move_file(folder / "CALENDARIO_2026-2.md", plan / "CALENDARIO_2026-2.md")
        )
        # PLAN DE CURSO docx from Entregas -> Plan curso
        for p in list(entregas.glob("PLAN DE CURSO*")) + list(entregas.glob("PLAN_DE_CURSO*")):
            log.append(move_file(p, plan / p.name))
        # NOTA operacional -> Plan curso (no es Acuerdo/Diagnóstico)
        for p in entregas.glob("NOTA*"):
            log.append(move_file(p, plan / p.name))

        # Report leftover in Entregas
        left = sorted(x.name for x in entregas.iterdir() if x.is_file())
        log.append(f"ENTREGAS {meta['folder']}: {left}")
    return log


def rows_literal(rows: list[list[str]], quote: str = '"') -> str:
    lines = []
    for r in rows:
        cells = ", ".join(f"{quote}{c}{quote}" for c in r)
        lines.append(f"            [{cells}],")
    return "\n".join(lines)


def patch_build(path: Path, rows: list[list[str]]) -> None:
    t = read_text(path)
    # Detect quote style around first cronograma row
    quote = '"' if '["1"' in t or "\"1\"" in t else "'"
    if "['1'" in t or '["1"' in t:
        pass
    # Prefer matching existing style in file
    if "['1'" in t or "['1'," in t:
        quote = "'"
    elif '["1"' in t:
        quote = '"'

    r1 = rows_literal(rows[:8], quote)
    r2 = rows_literal(rows[8:], quote)
    q = quote

    block = f'''table_content(
        prs, "Cronograma de clases (1–8)",
        [{q}#{q}, {q}Fecha{q}, {q}Tipo{q}, {q}Tema{q}],
        [
{r1}
        ],
        note={q}Presencialidad asistida: Clase 1 presencial · resto virtual · parciales presencial · festivos autónoma.{q},
        col_w=[0.7, 1.3, 1.5, 7.8],
        fs_body=11,
        idx=8,
    )

    table_content(
        prs, "Cronograma de clases (9–15)",
        [{q}#{q}, {q}Fecha{q}, {q}Tipo{q}, {q}Tema{q}],
        [
{r2}
        ],
        note={q}Ver Plan curso/PLAN_DE_CURSO_2026-2.md y Plan curso/CALENDARIO_2026-2.md.{q},
        col_w=[0.7, 1.3, 1.5, 7.8],
        fs_body=11,
        idx=9,
    )'''

    # Replace existing single cronograma table_content call
    pattern = re.compile(
        r"table_content\(\s*prs,\s*[\"']Cronograma de clases[^\"]*[\"'].*?"
        r"fs_body\s*=\s*11,\s*idx\s*=\s*\d+\s*,?\s*\)",
        re.S,
    )
    if not pattern.search(t):
        raise RuntimeError(f"No cronograma block in {path.name}")
    t = pattern.sub(block, t, count=1)

    # Bump idx after cronograma (old idx>=9 become +1), carefully by parsing calls
    # Simpler: replace known trailing idx values that follow cronograma
    # After insert, Proyecto Integrador was idx=9 -> 10, Recursos 10->11, Acuerdos 11->12, cierre 12->13
    def bump_idx(text: str) -> str:
        # Only bump idxs that appear AFTER the second cronograma (idx=9)
        marker = 'idx=9,\n    )'
        pos = text.find(marker)
        if pos < 0:
            marker = "idx=9,\n    )"
            pos = text.find(marker)
        if pos < 0:
            return text
        head, tail = text[: pos + len(marker)], text[pos + len(marker) :]

        def repl(m: re.Match) -> str:
            n = int(m.group(1))
            if n >= 9:
                return f"idx={n + 1}"
            return m.group(0)

        # bump idx=N in tail once each occurrence where N>=9
        # Do from high to low to avoid double bump
        for n in range(20, 8, -1):
            tail = re.sub(rf"\bidx\s*=\s*{n}\b", f"idx={n + 1}", tail)
        return head + tail

    t = bump_idx(t)

    # Ensure no Campus/listado pendientes remain
    lines = []
    for line in t.splitlines():
        low = line.lower()
        if "campus virtual:" in low and "pendiente" in low:
            continue
        if "listado de estudiantes" in low and "pendiente" in low:
            continue
        if "campus virtual uniajc:" in low and "pendiente" in low:
            continue
        if "meet" in low and "pendiente" in low and "url" in low:
            continue
        line = re.sub(r"\s*Listado:\s*\[PENDIENTE listado\]\.?", "", line)
        lines.append(line)
    t = "\n".join(lines)
    if not t.endswith("\n"):
        t += "\n"
    write_text(path, t)


def update_plan_refs() -> None:
    for meta in CURSOS.values():
        plan_md = ROOT / meta["folder"] / "Plan curso" / "PLAN_DE_CURSO_2026-2.md"
        if not plan_md.exists():
            continue
        t = read_text(plan_md)
        t = t.replace(
            "`Entregas docente/calendario_eventos_2026-2.csv`",
            "`Plan curso/calendario_eventos_2026-2.csv`",
        )
        t = t.replace(
            "`../CALENDARIO_2026-2.md`",
            "`CALENDARIO_2026-2.md`",
        )
        t = t.replace(
            "Calendario:** `../CALENDARIO_2026-2.md`",
            "Calendario:** `CALENDARIO_2026-2.md`",
        )
        # ensure calendario line points to Plan curso
        t = re.sub(
            r"- \*\*Calendario:\*\*.*",
            "- **Calendario:** `CALENDARIO_2026-2.md` · `.config/calendario/semestre_2026_2.json`",
            t,
        )
        write_text(plan_md, t if t.endswith("\n") else t + "\n")

        # NOTA files if moved
        for nota in (ROOT / meta["folder"] / "Plan curso").glob("NOTA*"):
            nt = read_text(nota)
            nt = nt.replace("Entregas docente/calendario_eventos", "Plan curso/calendario_eventos")
            nt = nt.replace("../CALENDARIO_2026-2.md", "CALENDARIO_2026-2.md")
            write_text(nota, nt if nt.endswith("\n") else nt + "\n")


def update_json_and_docs() -> None:
    jpath = ROOT / ".config/calendario/semestre_2026_2.json"
    data = json.loads(jpath.read_text(encoding="utf-8"))
    for key, curso in data["cursos"].items():
        folder = curso["folder"]
        curso["calendario"] = f"{folder}/Plan curso/CALENDARIO_2026-2.md"
        if "calendario_csv" not in curso:
            curso["calendario_csv"] = f"{folder}/Plan curso/calendario_eventos_2026-2.csv"
        else:
            curso["calendario_csv"] = f"{folder}/Plan curso/calendario_eventos_2026-2.csv"
    jpath.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    readme = ROOT / ".config/calendario/README_eventos_csv.md"
    write_text(
        readme,
        "# Calendario de eventos CSV 2026-2\n\n"
        "Archivos `eventos_*_2026-2.csv` (copia en `.config/calendario/`) y "
        "`<Curso>/Plan curso/calendario_eventos_2026-2.csv`: 15 filas/clase por curso, UTF-8 con BOM.\n\n"
        "Cuando tengas el listado de estudiantes, importa el CSV a Excel/Google Sheets o genera invitaciones "
        "(una fila = una clase; `es_parcial=si` marca parciales síncronos; "
        "`tipo_clase` = `presencial` | `virtual` | `autonoma`).\n\n"
        "Modalidad por sesión: Clase 1 presencial; resto virtual; parciales presencial; festivos autónoma.\n",
    )

    # generar_semestre output path
    gen = ROOT / ".config/calendario/generar_semestre_2026_2.py"
    gt = read_text(gen)
    gt = gt.replace(
        'cal_path = folder / "CALENDARIO_2026-2.md"',
        'cal_path = folder / "Plan curso" / "CALENDARIO_2026-2.md"',
    )
    gt = gt.replace(
        "[Detalle en CALENDARIO_2026-2.md del curso.]",
        "[Detalle en Plan curso/CALENDARIO_2026-2.md.]",
    )
    write_text(gen, gt if gt.endswith("\n") else gt + "\n")


def update_rules_agents() -> None:
    rule = ROOT / ".cursor/rules/uniajc-docente.mdc"
    t = read_text(rule)
    t = t.replace(
        "Fuente de calendario/cortes: `config/calendario/semestre_2026_2.json` (+ `CALENDARIO_2026-2.md` por curso).",
        "Fuente de calendario/cortes: `config/calendario/semestre_2026_2.json` "
        "(+ `Plan curso/CALENDARIO_2026-2.md` y `Plan curso/calendario_eventos_2026-2.csv` por curso).",
    )
    t = t.replace(
        "  - Cierre: día/hora semanal; cronograma desde `CALENDARIO_2026-2.md`.",
        "  - Cierre: día/hora semanal; cronograma desde `Plan curso/CALENDARIO_2026-2.md`.",
    )
    old_struct = (
        "  Plan curso/          Microcurrículo + Plan_de_curso\n"
        "  Entregas docente/    Acuerdo pedagógico, Diagnóstico, PLAN DE CURSO\n"
    )
    new_struct = (
        "  Plan curso/          Microcurrículo + Plan_de_curso + PLAN_DE_CURSO_2026-2.md + "
        "CALENDARIO_2026-2.md + calendario_eventos_2026-2.csv\n"
        "  Entregas docente/    Acuerdo pedagógico + Diagnóstico (sin calendarios ni planes)\n"
    )
    if old_struct in t:
        t = t.replace(old_struct, new_struct)
    else:
        t = t.replace(
            "Entregas docente/    Acuerdo pedagógico, Diagnóstico, PLAN DE CURSO",
            "Entregas docente/    Acuerdo pedagógico + Diagnóstico (sin calendarios ni planes)",
        )
        t = t.replace(
            "Plan curso/          Microcurrículo + Plan_de_curso",
            "Plan curso/          Microcurrículo + Plan_de_curso + CALENDARIO + CSV eventos + PLAN_DE_CURSO_2026-2.md",
        )
    t = t.replace(
        "  CALENDARIO_2026-2.md                     ← calendario del periodo\n",
        "  Plan curso/CALENDARIO_2026-2.md           ← calendario del periodo\n",
    )
    write_text(rule, t if t.endswith("\n") else t + "\n")

    for ap in [
        ROOT / ".claude/agents/disenador-curricular-uniajc.md",
        ROOT / ".cursor/agents/disenador-curricular-uniajc.md",
    ]:
        if not ap.exists():
            continue
        at = read_text(ap)
        at = at.replace(
            "Microcurrículo y/o Plan de curso de la asignatura (en `Plan curso/` o `Entregas docente/`)",
            "Microcurrículo y/o Plan de curso de la asignatura (en `Plan curso/`)",
        )
        at = at.replace(
            "a) `Entregas docente/PLAN DE CURSO…` y `ACUERDO PEDAGOGICO…` (evaluación + metodología del grupo)\n"
            "   - b) `Plan curso/Plan_de_curso…` + `Microcurrriculo…`",
            "a) `Entregas docente/ACUERDO PEDAGOGICO…` (evaluación + metodología del grupo)\n"
            "   - b) `Plan curso/` (Plan_de_curso, Microcurrículo, PLAN_DE_CURSO_2026-2.md, CALENDARIO, CSV eventos)",
        )
        # structure block if present
        at = at.replace(
            "<Curso>/Entregas docente/  (Acuerdo pedagógico, Diagnóstico, PLAN DE CURSO)",
            "<Curso>/Entregas docente/  (Acuerdo pedagógico, Diagnóstico)",
        )
        at = at.replace(
            "<Curso>/Plan curso/  (Microcurrículo + Plan_de_curso)",
            "<Curso>/Plan curso/  (Microcurrículo + Plan_de_curso + PLAN_DE_CURSO_2026-2 + CALENDARIO + CSV)",
        )
        # cronograma note in entregable
        if "2 slides" not in at and "Cronograma de clases" in at:
            at = at.replace(
                "8. Cronograma de clases (tabla del Plan de curso 2026-2)",
                "8. Cronograma de clases (tablas 1–8 y 9–15; todas las clases visibles, sin resumen en pie)",
            )
        write_text(ap, at if at.endswith("\n") else at + "\n")

    for dp in [
        ROOT / ".claude/agents/uniajc-dudas-material.md",
        ROOT / ".cursor/agents/uniajc-dudas-material.md",
    ]:
        if not dp.exists():
            continue
        dt = read_text(dp)
        dt = dt.replace(
            "(+ `CALENDARIO_2026-2.md` por curso)",
            "(+ `Plan curso/CALENDARIO_2026-2.md` por curso)",
        )
        dt = dt.replace(
            "`Plan curso/`, `Entregas docente/`, `Clases/`, `Kit docente/`, `CALENDARIO_2026-2.md`",
            "`Plan curso/` (plan+calendario+CSV), `Entregas docente/` (Acuerdo+Diagnóstico), `Clases/`, `Kit docente/`",
        )
        write_text(dp, dt if dt.endswith("\n") else dt + "\n")

    uj = ROOT / ".config/universidades/uniajc.json"
    u = json.loads(uj.read_text(encoding="utf-8"))
    u["estructura_carpetas_curso"]["esquema_existente"] = [
        "<Curso>/Plan curso/  (Microcurrículo + Plan_de_curso + PLAN_DE_CURSO_2026-2.md + CALENDARIO_2026-2.md + calendario_eventos_2026-2.csv)",
        "<Curso>/Entregas docente/  (Acuerdo pedagógico, Diagnóstico)",
        "<Curso>/Kit docente/Clase N/  (Guion, Quiz, Codigo — material docente)",
        "<Curso>/Clases/Clase N/  (presentación estudiante + taller)",
        "<Curso>/Parciales/",
        "<Curso>/Clases grabadas/",
    ]
    u["fuente_de_temas"]["orden"] = [
        "Plan de curso / PLAN_DE_CURSO_2026-2.md en Plan curso/",
        "Microcurrículo oficial (.docx) — saberes, RAAs, estrategia didáctica",
        "Acuerdo pedagógico del grupo en Entregas docente/ (evaluación, metodología)",
        "Cronograma Plan curso/CALENDARIO_2026-2.md + CSV eventos",
        "Cronograma ya materializado en Clases/Clase N - <Tema> (si existe, respétalo)",
    ]
    uj.write_text(json.dumps(u, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def rebuild() -> None:
    for meta in CURSOS.values():
        script = ROOT / ".config/slides" / meta["build"]
        print("BUILD", script.name)
        r = subprocess.run(
            [sys.executable, str(script)], cwd=str(ROOT), capture_output=True, text=True
        )
        print(r.stdout)
        if r.returncode != 0:
            print(r.stderr)
            raise SystemExit(f"Falló {script.name}")


def verify_pptx() -> None:
    from pptx import Presentation

    for meta in CURSOS.values():
        folder = ROOT / meta["folder"] / "Clases"
        pptxs = list(folder.glob("Presentacion del Curso*.pptx"))
        if not pptxs:
            print("MISSING pptx", meta["folder"])
            continue
        prs = Presentation(str(pptxs[0]))
        found_nums = set()
        bad = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cell0 = row.cells[0].text.strip()
                        if cell0.isdigit():
                            found_nums.add(int(cell0))
                        texts.append(" | ".join(c.text for c in row.cells))
            blob = "\n".join(texts)
            if "Clases 9" in blob and "–" in blob:
                bad.append("summary note")
            if "PENDIENTE listado" in blob or (
                "URL Campus Virtual" in blob and "pendiente" in blob.lower()
            ):
                bad.append("placeholder")
        print(
            pptxs[0].name,
            "nums",
            sorted(found_nums),
            "complete" if found_nums >= set(range(1, 16)) else "INCOMPLETE",
            "bad",
            bad,
        )


def main() -> None:
    print("=== REORG ===")
    for line in reorganize():
        print(line)
    print("=== PLAN REFS ===")
    update_plan_refs()
    print("=== JSON/DOCS ===")
    update_json_and_docs()
    print("=== RULES/AGENTS ===")
    update_rules_agents()
    print("=== BUILDS ===")
    for key, meta in CURSOS.items():
        patch_build(ROOT / ".config/slides" / meta["build"], meta["rows"])
        print("patched", meta["build"])
    print("=== REBUILD ===")
    rebuild()
    print("=== VERIFY ===")
    verify_pptx()
    print("DONE")


if __name__ == "__main__":
    main()
